#!/usr/bin/env python3
"""
投行部智能图片生成器 v4.0 — 精致桌面应用
基于 customtkinter 现代化 UI 框架 · 云雾API (yunwu.ai) · OpenAI 兼容协议

三大功能页:
  · 配置      —— API Key / 接口地址 / 模型管理(查询·手动·多保存) / 画质 / 默认保存路径
  · 文本对话  —— 多轮对话(展示思考过程) / 上传任意格式文件分析 / 回复可复制·查看全文
  · 单次出图  —— 单条提示词，一次可出多张
  · 批量出图  —— 提示词列表(默认5条可增减)，每条单独出 1 张，最多 30 张

协议: OpenAI 兼容同步协议
  · 文生图   —— POST /v1/images/generations (JSON)
  · 图生图   —— POST /v1/images/edits (multipart 文件直传)
  · 查询模型 —— GET  /v1/models
特性: 自定义模型管理 / 文生图+图生图 / 并发出图(最多30) / 参考图上传 / 调用日志
打包: pyinstaller --onefile --windowed --name="投行部智能图片生成器" image_generator_app.py
"""

import customtkinter as ctk
from tkinter import filedialog, messagebox
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
import requests
import json
import os
import base64
import re
from datetime import datetime
from pathlib import Path
from PIL import Image, ImageTk
from io import BytesIO

# ============================================================
# 全局配置
# ============================================================
APP_NAME = "投行部智能图片生成器"
APP_VERSION = "v4.5"

# 批量并发上限（云雾 API 支持高并发，最多 30 路同时出图）
MAX_CONCURRENCY = 30
BATCH_MAX = 30    # 批量出图：最多提示词条数（每条出 1 张）

# gpt-image-2 画质档位（quality 字段），默认 auto
QUALITY_OPTIONS = ["auto", "low", "medium", "high"]

# 模型选择弹窗：保存上限 + 单次渲染上限（防止数百模型一次性渲染卡顿）
MODEL_SAVE_LIMIT = 10
MODEL_RENDER_LIMIT = 80

# 统一风格指令：批量出图勾选「锁定参考图风格」时，自动拼到每条提示词前，
# 强制模型把参考图当作视觉模板（配色/版式/字体/质感保持一致），只替换内容。
STYLE_LOCK_DIRECTIVE = (
    "严格沿用所附参考图的视觉风格：保持完全一致的整体配色方案、版式布局结构、"
    "字体风格、图标与装饰元素的设计语言、画面质感与光影氛围。"
    "参考图作为风格模板，仅按下方描述替换具体内容文字与主体，"
    "其余视觉风格元素必须与参考图统一，确保系列图片风格连贯。具体内容："
)
BATCH_DEFAULT_ROWS = 5  # 批量页默认提示词行数

# 文本/分析模型预设（OpenAI 协议 /v1/chat/completions）。
# 名称含 REASONING_HINTS 的模型会尝试展示思考过程（reasoning_content）。
CHAT_MODELS = [
    "gpt-4o",
    "gpt-4o-mini",
    "gpt-4.1",
    "o1",
    "o3-mini",
    "deepseek-chat",
    "deepseek-reasoner",
    "claude-3-7-sonnet-20250219",
]
REASONING_HINTS = ("o1", "o3", "reasoner", "r1", "thinking", "qwq", "deep-think", "think")

# 上传文件分析：可直接读取为文本的扩展名 + 单文件读取上限
TEXT_READABLE_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".xml", ".yaml", ".yml",
    ".log", ".ini", ".conf", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h",
    ".cs", ".go", ".rs", ".rb", ".php", ".html", ".htm", ".css", ".sql", ".sh",
    ".bat", ".vue", ".jsx", ".tsx",
}
# 需要专用库解析的二进制文档格式（不在 TEXT_READABLE_EXTS 中，由 _build_files_context 单独处理）
BINARY_DOC_EXTS = {".docx", ".xlsx", ".pptx", ".pdf"}
FILE_READ_MAX = 200 * 1024  # 单文件最多读取 200KB 文本，避免 prompt 过长

# 色彩方案
COLORS = {
    "bg": "#0F1117",
    "card": "#1A1D27",
    "card_border": "#2A2D37",
    "text_primary": "#E8EAED",
    "text_secondary": "#9AA0A6",
    "accent": "#4F8FFF",
    "accent_hover": "#6BA4FF",
    "success": "#34A853",
    "warning": "#FBBC04",
    "error": "#EA4335",
    "surface": "#242835",
    "divider": "#3C4043",
    "log_info": "#4F8FFF",
    "log_success": "#34A853",
    "log_error": "#EA4335",
    "log_request": "#FBBC04",
}

# 尺寸字典：键为下拉显示文案，值为云雾 gpt-image-2 文档支持的 size 取值。
# 文档约束：最长边≤3840px、宽高均为16的倍数、长短边比≤3:1、总像素 655360~8294400。
# 详见云雾 apifox「创建 gpt-image-2」接口文档。
SIZE_MAP = {
    "自动（由模型决定）":        {"pixel": "auto",      "ratio": "auto"},
    "1:1  正方形 1024×1024":     {"pixel": "1024x1024", "ratio": "1:1"},
    "3:2  横版 1536×1024":       {"pixel": "1536x1024", "ratio": "3:2"},
    "2:3  竖版 1024×1536":       {"pixel": "1024x1536", "ratio": "2:3"},
    "1:1  2K 正方形 2048×2048":  {"pixel": "2048x2048", "ratio": "1:1"},
    "16:9 2K 横版 2048×1152":    {"pixel": "2048x1152", "ratio": "16:9"},
    "16:9 4K 横版 3840×2160":    {"pixel": "3840x2160", "ratio": "16:9"},
    "9:16 4K 竖版 2160×3840":    {"pixel": "2160x3840", "ratio": "9:16"},
}
DEFAULT_SIZE_KEY = "1:1  正方形 1024×1024"

CONFIG_FILE = Path.home() / ".gpt_image_gen_config.json"
DEFAULT_CONFIG = {
    "api_key": "",
    "api_base": "https://yunwu.ai/v1",
    "save_dir": str(Path.home() / "Pictures"),
    "filename_prefix": "ai_image",
    "last_size_key": DEFAULT_SIZE_KEY,
    # 用户保存的模型列表（云雾 API 内的模型 ID），下拉框从此读取
    "models": ["gpt-image-2"],
    "model": "gpt-image-2",
    "quality": "auto",
    # 智能分析（文本对话）默认入参
    "chat_model": "gpt-4o",
    "chat_system": "你是投行部的专业分析助手，回答严谨、简洁、有条理。",
    "chat_temperature": 0.7,
    "chat_max_tokens": 2048,
    "chat_stream": True,
}


def is_reasoning_model(name: str) -> bool:
    """模型名包含推理关键字时，认为是推理模型（展示思考过程）。"""
    n = (name or "").lower()
    return any(h in n for h in REASONING_HINTS)


def load_config():
    if CONFIG_FILE.exists():
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                return {**DEFAULT_CONFIG, **json.load(f)}
        except Exception:
            pass
    return DEFAULT_CONFIG.copy()


def save_config(cfg: dict):
    try:
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump(cfg, f, ensure_ascii=False, indent=2)
    except Exception:
        pass


# ============================================================
# API 客户端 — OpenAI 兼容协议（云雾 yunwu.ai）
# ============================================================
def _normalize_openai_base(base: str) -> str:
    """把用户填的各种形态规整成以 /v1 结尾的根。
    接受 https://x、https://x/、https://x/v1、https://x/v1/images/generations 等。"""
    b = (base or "").strip().rstrip("/")
    if not b:
        return "https://yunwu.ai/v1"
    # 砍掉常见的接口尾巴
    for tail in ("/images/generations", "/images/edits",
                 "/chat/completions", "/completions"):
        if b.endswith(tail):
            b = b[: -len(tail)]
            break
    b = b.rstrip("/")
    if not b.endswith("/v1"):
        b = b + "/v1"
    return b


def _safe_json(r, context: str):
    if not r.text or not r.text.strip():
        raise RuntimeError(f"{context}: API 返回空响应 (HTTP {r.status_code})")
    try:
        return r.json()
    except Exception:
        raise RuntimeError(f"{context}: API 返回非 JSON (HTTP {r.status_code})\n{r.text[:300]}")


def _extract_error(data, status_code):
    """从 OpenAI 风格错误体中提取人类可读消息。"""
    err = data.get("error") if isinstance(data, dict) else None
    if isinstance(err, dict):
        return err.get("message") or err.get("code") or f"HTTP {status_code}"
    if isinstance(err, str):
        return err
    return data.get("message", f"HTTP {status_code}") if isinstance(data, dict) else f"HTTP {status_code}"


def _decode_image_item(item):
    """从返回的 data[i] 中拿到图片字节：优先 b64_json，否则下载 url。"""
    if item.get("b64_json"):
        return base64.b64decode(item["b64_json"])
    if item.get("url"):
        dl = requests.get(item["url"], timeout=120)
        if dl.status_code != 200:
            raise RuntimeError(f"图片下载失败 (HTTP {dl.status_code})")
        return dl.content
    raise RuntimeError("返回结果无 b64_json / url")


def openai_list_models(api_base, api_key, timeout=20):
    """查询云雾 API 可用模型列表：GET /v1/models。
    返回模型 ID 字符串列表（已排序）。失败抛 RuntimeError。"""
    base = _normalize_openai_base(api_base)
    url = f"{base}/models"
    headers = {"Authorization": f"Bearer {api_key}"}
    r = requests.get(url, headers=headers, timeout=timeout)
    data = _safe_json(r, "查询模型")
    if r.status_code >= 400 or (isinstance(data, dict) and "error" in data):
        raise RuntimeError(_extract_error(data, r.status_code))
    items = data.get("data", []) if isinstance(data, dict) else []
    ids = [it.get("id") for it in items if isinstance(it, dict) and it.get("id")]
    return sorted(set(ids)), {"url": url, "count": len(ids)}


def openai_generations(api_base, api_key, model, prompt, size_pixel,
                       quality="auto", n=1, timeout=300):
    """文生图：POST /v1/images/generations (application/json)。
    返回 (image_bytes, log_payload)。"""
    base = _normalize_openai_base(api_base)
    url = f"{base}/images/generations"
    headers = {"Authorization": f"Bearer {api_key}",
               "Content-Type": "application/json"}
    payload = {"model": model, "prompt": prompt, "n": n, "size": size_pixel}
    if quality and quality != "auto":
        payload["quality"] = quality

    r = requests.post(url, json=payload, headers=headers, timeout=timeout)
    data = _safe_json(r, "文生图")
    if r.status_code >= 400 or (isinstance(data, dict) and "error" in data):
        raise RuntimeError(_extract_error(data, r.status_code))
    items = data.get("data", [])
    if not items:
        raise RuntimeError("文生图返回空结果")
    log = {"url": url, "payload": {**payload}}
    return _decode_image_item(items[0]), log


def openai_edits(api_base, api_key, model, prompt, size_pixel,
                 ref_files, quality="auto", n=1, timeout=300):
    """图生图：POST /v1/images/edits (multipart/form-data 文件直传)。
    ref_files: [(filename, bytes, mime), ...]，最多 16 张，用 image[] 字段上传。
    返回 (image_bytes, log_payload)。"""
    base = _normalize_openai_base(api_base)
    url = f"{base}/images/edits"
    headers = {"Authorization": f"Bearer {api_key}"}  # multipart 边界由 requests 自动设置

    data_fields = {"model": (None, model), "prompt": (None, prompt),
                   "n": (None, str(n)), "size": (None, size_pixel)}
    if quality and quality != "auto":
        data_fields["quality"] = (None, quality)

    files = []
    for fname, fbytes, mime in (ref_files or [])[:16]:
        files.append(("image[]", (fname, fbytes, mime)))
    # 合并文本字段与文件字段一起作为 multipart
    multipart = list(data_fields.items()) + files

    r = requests.post(url, files=multipart, headers=headers, timeout=timeout)
    data = _safe_json(r, "图生图")
    if r.status_code >= 400 or (isinstance(data, dict) and "error" in data):
        raise RuntimeError(_extract_error(data, r.status_code))
    items = data.get("data", [])
    if not items:
        raise RuntimeError("图生图返回空结果")
    log = {"url": url,
           "payload": {"model": model, "prompt": prompt[:80], "n": n,
                       "size": size_pixel, "image": f"[{len(files)} 张参考图]"}}
    return _decode_image_item(items[0]), log


def openai_chat_stream(openai_base, openai_key, model, messages,
                       temperature=0.7, max_tokens=2048,
                       on_reasoning=None, on_content=None, stop_flag=None):
    """OpenAI 协议流式对话 /v1/chat/completions。
    on_reasoning(text_delta): 推理增量回调（如有 reasoning_content）
    on_content(text_delta):   正文增量回调
    stop_flag(): 返回 True 则中断
    返回 (full_reasoning, full_content)。"""
    base = _normalize_openai_base(openai_base)
    url = f"{base}/chat/completions"
    headers = {"Authorization": f"Bearer {openai_key}",
               "Content-Type": "application/json",
               "Accept": "text/event-stream"}
    payload = {"model": model, "messages": messages,
               "temperature": temperature, "stream": True}
    if max_tokens:
        payload["max_tokens"] = max_tokens

    full_reason, full_content = [], []
    with requests.post(url, json=payload, headers=headers,
                       stream=True, timeout=300) as r:
        if r.status_code >= 400:
            # 非流式错误：直接读全文
            raise RuntimeError(f"对话请求失败 (HTTP {r.status_code}): {r.text[:300]}")
        # 逐小块读取 SSE 流，避免 iter_lines 被换行/代理缓冲阻塞
        sse_buf = ""
        for raw_bytes in r.iter_content(chunk_size=None):
            if stop_flag and stop_flag():
                break
            if not raw_bytes:
                continue
            # 强制 UTF-8 解码（SSE 流可能无 charset 声明）
            try:
                text = raw_bytes.decode("utf-8")
            except Exception:
                text = raw_bytes.decode("utf-8", errors="replace")
            sse_buf += text
            # 按 \n\n 分割完整的 SSE 事件（每个事件以空行结尾）
            while "\n\n" in sse_buf:
                event_text, sse_buf = sse_buf.split("\n\n", 1)
                # 解析事件内的每行
                for line in event_text.split("\n"):
                    line = line.strip()
                    if line.startswith("data:"):
                        line = line[5:].strip()
                    elif line.startswith(":"):
                        # SSE 注释行，跳过
                        continue
                    else:
                        continue
                    if line == "[DONE]":
                        return "".join(full_reason), "".join(full_content)
                    try:
                        chunk = json.loads(line)
                    except Exception:
                        continue
                    choices = chunk.get("choices") or []
                    if not choices:
                        continue
                    delta = choices[0].get("delta", {}) or {}
                    # 兼容多种推理字段名：reasoning_content / reasoning / think / thought
                    rc = (delta.get("reasoning_content") or delta.get("reasoning")
                          or delta.get("think") or delta.get("thought"))
                    if rc:
                        full_reason.append(rc)
                        if on_reasoning:
                            on_reasoning(rc)
                    c = delta.get("content")
                    if c:
                        full_content.append(c)
                        if on_content:
                            on_content(c)
        # 处理缓冲区中可能残留的数据
        if sse_buf.strip():
            for line in sse_buf.split("\n"):
                line = line.strip()
                if line.startswith("data:"):
                    line = line[5:].strip()
                else:
                    continue
                if line == "[DONE]":
                    break
                try:
                    chunk = json.loads(line)
                except Exception:
                    continue
                choices = chunk.get("choices") or []
                if not choices:
                    continue
                delta = choices[0].get("delta", {}) or {}
                rc = (delta.get("reasoning_content") or delta.get("reasoning")
                      or delta.get("think") or delta.get("thought"))
                if rc:
                    full_reason.append(rc)
                    if on_reasoning:
                        on_reasoning(rc)
                c = delta.get("content")
                if c:
                    full_content.append(c)
                    if on_content:
                        on_content(c)
    return "".join(full_reason), "".join(full_content)


def openai_chat_once(openai_base, openai_key, model, messages,
                     temperature=0.7, max_tokens=2048):
    """非流式对话，返回 (reasoning, content)。用于关闭流式时。"""
    base = _normalize_openai_base(openai_base)
    url = f"{base}/chat/completions"
    headers = {"Authorization": f"Bearer {openai_key}",
               "Content-Type": "application/json"}
    payload = {"model": model, "messages": messages,
               "temperature": temperature, "stream": False}
    if max_tokens:
        payload["max_tokens"] = max_tokens
    r = requests.post(url, json=payload, headers=headers, timeout=300)
    data = _safe_json(r, "对话")
    if r.status_code >= 400 or (isinstance(data, dict) and "error" in data):
        raise RuntimeError(_extract_error(data, r.status_code))
    msg = (data.get("choices") or [{}])[0].get("message", {}) or {}
    return msg.get("reasoning_content", "") or "", msg.get("content", "") or ""




# ============================================================
# 页面上下文：持有某个生成页的 UI 组件引用（一览/进度/按钮）
# ============================================================
class PageContext:
    def __init__(self):
        self.gallery = None
        self.gallery_photos = []
        self.gallery_col = 0
        self.gallery_row = 0
        self.placeholder = None
        self.progress_frame = None
        self.progress_bar = None
        self.lbl_progress = None
        self.btn_generate = None
        # 参考图相关（单次/批量各自独立一套）
        self.ref_paths = []
        self.ref_photos = []
        self.ref_thumb_frame = None
        self.lbl_ref_count = None
        self.entry_ref_urls = None
        self.style_lock_var = None    # 批量页"锁定参考图风格"开关


# ============================================================
# 主窗口
# ============================================================
class ImageGeneratorApp:
    def __init__(self):
        self.config = load_config()
        ctk.set_appearance_mode("dark")
        ctk.set_default_color_theme("blue")

        self.win = ctk.CTk()
        self.win.title(APP_NAME)
        self.win.geometry("760x840")
        self.win.minsize(620, 640)
        self.win.configure(fg_color=COLORS["bg"])

        self.win.update_idletasks()
        sw, sh = self.win.winfo_screenwidth(), self.win.winfo_screenheight()
        self.win.geometry(f"760x840+{(sw-760)//2}+{(sh-840)//2}")

        self.generating = False
        self.batch_rows = []          # 批量页的提示词行 [{frame, entry, idx_lbl}]
        self.log_expanded = False
        self.MAX_LOG_LINES = 500

        # 两个生成页各自的上下文
        self.ctx_single = PageContext()
        self.ctx_batch = PageContext()

        # 智能分析（文本对话）状态
        self.chat_history = []            # [{"role","content"}] 不含 system
        self.chat_streaming = False
        self.chat_stop = False
        self.chat_placeholder = None
        self.chat_files = []              # 待分析文件 [{"path","name","ext"}]
        self._chat_cur_content_lbl = None
        self._chat_cur_reason_box = None
        self._chat_cur_content_buf = ""
        self._chat_cur_reason_buf = ""

        self._build()

    # ================================================================
    #  UI 工具方法
    # ================================================================
    def _mk_card(self, parent, **kw):
        return ctk.CTkFrame(parent, fg_color=COLORS["card"],
                            border_color=COLORS["card_border"],
                            border_width=1, corner_radius=12, **kw)

    def _mk_label(self, parent, text, size=13, color=None, weight="normal", **kw):
        ff = "Microsoft YaHei UI" if os.name == "nt" else "Segoe UI"
        return ctk.CTkLabel(parent, text=text, font=(ff, size, weight),
                            text_color=color or COLORS["text_primary"], **kw)

    def _mk_entry(self, parent, placeholder="", **kw):
        return ctk.CTkEntry(parent, height=38, fg_color=COLORS["surface"],
                            border_color=COLORS["divider"], border_width=1,
                            corner_radius=8, text_color=COLORS["text_primary"],
                            placeholder_text=placeholder,
                            placeholder_text_color=COLORS["text_secondary"], **kw)

    def _mk_btn(self, parent, text, style="primary", **kw):
        if style == "primary":
            return ctk.CTkButton(parent, text=text, height=42, corner_radius=10,
                                 fg_color=COLORS["accent"],
                                 hover_color=COLORS["accent_hover"],
                                 font=("Microsoft YaHei UI", 14, "bold"), **kw)
        elif style == "secondary":
            return ctk.CTkButton(parent, text=text, height=34, corner_radius=8,
                                 fg_color=COLORS["surface"],
                                 hover_color=COLORS["card_border"],
                                 border_color=COLORS["divider"], border_width=1,
                                 font=("Microsoft YaHei UI", 12), **kw)
        else:
            return ctk.CTkButton(parent, text=text, height=34, corner_radius=8,
                                 fg_color="transparent", hover_color=COLORS["surface"],
                                 font=("Microsoft YaHei UI", 11),
                                 text_color=COLORS["text_secondary"], **kw)

    def _mk_optionmenu(self, parent, values, width=170, height=34, command=None):
        return ctk.CTkOptionMenu(
            parent, values=values, height=height, width=width,
            fg_color=COLORS["surface"], button_color=COLORS["accent"],
            button_hover_color=COLORS["accent_hover"],
            text_color=COLORS["text_primary"],
            font=("Microsoft YaHei UI", 12), corner_radius=8,
            command=command)

    def _bind_tooltip(self, widget, text):
        """给控件绑定悬停提示（轻量级 Toplevel 实现）。"""
        tip = {"win": None}

        def show(_event=None):
            if tip["win"] is not None:
                return
            x = widget.winfo_rootx() + 20
            y = widget.winfo_rooty() + widget.winfo_height() + 6
            tw = ctk.CTkToplevel(self.win)
            tw.overrideredirect(True)
            tw.attributes("-topmost", True)
            tw.geometry(f"+{x}+{y}")
            tw.configure(fg_color=COLORS["card"])
            lbl = ctk.CTkLabel(tw, text=text, font=("Microsoft YaHei UI", 11),
                               text_color=COLORS["text_primary"],
                               fg_color=COLORS["card"], wraplength=360,
                               justify="left")
            lbl.pack(padx=10, pady=6)
            tip["win"] = tw

        def hide(_event=None):
            if tip["win"] is not None:
                tip["win"].destroy()
                tip["win"] = None

        widget.bind("<Enter>", show)
        widget.bind("<Leave>", hide)

    def _info_dot(self, parent, text):
        """生成一个精致的圆形 ⓘ 提示标记，鼠标悬停显示参数定义。
        返回该标记控件（调用方自行 pack/grid）。"""
        dot = ctk.CTkLabel(
            parent, text="ⓘ", width=18, height=18,
            font=("Microsoft YaHei UI", 12),
            text_color=COLORS["text_secondary"],
            fg_color="transparent")
        # 悬停时变色，强化「可交互」的视觉反馈
        def _enter(_e=None):
            dot.configure(text_color=COLORS["accent"])
        def _leave(_e=None):
            dot.configure(text_color=COLORS["text_secondary"])
        dot.bind("<Enter>", _enter)
        dot.bind("<Leave>", _leave)
        self._bind_tooltip(dot, text)
        return dot

    # ================================================================
    #  整体布局
    # ================================================================
    def _build(self):
        # --- 顶部标题 ---
        header = ctk.CTkFrame(self.win, fg_color="transparent", height=50)
        header.pack(fill="x", padx=24, pady=(18, 4))
        header.pack_propagate(False)

        tf = ctk.CTkFrame(header, fg_color="transparent")
        tf.pack(side="left")
        self._mk_label(tf, "🎨", size=24).pack(side="left", padx=(0, 8))
        self._mk_label(tf, "投行部", size=20, weight="bold",
                       color=COLORS["accent"]).pack(side="left")
        self._mk_label(tf, "智能图片生成器", size=13,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(6, 0))

        self._mk_label(header, APP_VERSION, size=10,
                       color=COLORS["text_secondary"]).pack(side="right", pady=(10, 0))

        # --- 导航（三 Tab）---
        nav = ctk.CTkFrame(self.win, fg_color="transparent")
        nav.pack(fill="x", padx=24, pady=(2, 4))

        self.nav_config = self._mk_btn(nav, "⚙️  配置", "ghost",
                                       command=lambda: self._switch_tab("config"))
        self.nav_config.pack(side="left", padx=(0, 4))
        self.nav_analyze = self._mk_btn(nav, "💬  文本对话", "ghost",
                                        command=lambda: self._switch_tab("analyze"))
        self.nav_analyze.pack(side="left", padx=4)
        self.nav_single = self._mk_btn(nav, "🖼️  单次出图", "ghost",
                                       command=lambda: self._switch_tab("single"))
        self.nav_single.pack(side="left", padx=4)
        self.nav_batch = self._mk_btn(nav, "🗂️  批量出图", "ghost",
                                      command=lambda: self._switch_tab("batch"))
        self.nav_batch.pack(side="left", padx=4)

        # --- 主内容区 ---
        self.main = ctk.CTkFrame(self.win, fg_color="transparent")
        self.main.pack(fill="both", expand=True, padx=24)

        self.page_config = self._build_config_page()
        self.page_analyze = self._build_analyze_page()
        self.page_single = self._build_single_page()
        self.page_batch = self._build_batch_page()

        self.page_config.pack_forget()
        self.page_analyze.pack_forget()
        self.page_batch.pack_forget()
        self.page_single.pack(fill="both", expand=True, pady=(4, 0))
        self.current_page = "single"
        self._highlight_nav("single")

        # --- 日志面板 (折叠) ---
        self._build_log_panel()

        # --- 底栏 ---
        footer = ctk.CTkFrame(self.win, fg_color="transparent", height=36)
        footer.pack(fill="x", padx=24, pady=(2, 10))
        footer.pack_propagate(False)

        self.status_bar = self._mk_label(footer, "就绪，输入提示词后点击「开始生成」",
                                         size=11, color=COLORS["text_secondary"])
        self.status_bar.pack(side="left")

        self.lbl_api_status = self._mk_label(footer, "🔴 未配置 API", size=11,
                                             color=COLORS["error"])
        self.lbl_api_status.pack(side="right")
        self._check_api_status()

    # ================================================================
    #  日志面板
    # ================================================================
    def _build_log_panel(self):
        self.log_container = ctk.CTkFrame(self.win, fg_color="transparent")
        self.log_container.pack(fill="x", padx=24, pady=(0, 0))

        toggle_bar = ctk.CTkFrame(self.log_container, fg_color=COLORS["card"],
                                  corner_radius=10, height=32)
        toggle_bar.pack(fill="x")
        toggle_bar.pack_propagate(False)

        self.btn_log_toggle = self._mk_btn(toggle_bar, "📋  调用日志  ▸", "ghost",
                                           command=self._toggle_log)
        self.btn_log_toggle.pack(side="left", padx=(12, 0))

        self.lbl_log_hint = self._mk_label(toggle_bar, "点击展开查看 API 调用详情",
                                           size=10, color=COLORS["text_secondary"])
        self.lbl_log_hint.pack(side="left", padx=(10, 0))

        btn_clear_log = self._mk_btn(toggle_bar, "清空", "ghost",
                                     command=self._clear_log)
        btn_clear_log.pack(side="right", padx=(0, 8), pady=4)

        self.log_frame = ctk.CTkFrame(self.log_container, fg_color=COLORS["card"],
                                      corner_radius=0)
        self.log_text = ctk.CTkTextbox(
            self.log_frame, height=140,
            fg_color=COLORS["surface"], border_color=COLORS["divider"],
            border_width=1, corner_radius=6, wrap="word",
            font=("Consolas", 10), text_color=COLORS["text_secondary"],
            state="disabled",
        )
        self.log_text.pack(fill="both", expand=True, padx=8, pady=(6, 8))

    def _toggle_log(self):
        if self.log_expanded:
            self.log_frame.pack_forget()
            self.btn_log_toggle.configure(text="📋  调用日志  ▸")
            self.lbl_log_hint.configure(text="点击展开查看 API 调用详情")
            self.log_expanded = False
        else:
            self.log_frame.pack(fill="x", before=self.log_container.winfo_children()[0])
            self.btn_log_toggle.configure(text="📋  调用日志  ▾")
            self.lbl_log_hint.configure(text="")
            self.log_expanded = True
            self._scroll_log_to_bottom()

    def _clear_log(self):
        self.log_text.configure(state="normal")
        self.log_text.delete("1.0", "end")
        self.log_text.configure(state="disabled")

    def _append_log(self, text, tag=None):
        def _do():
            ts = datetime.now().strftime("%H:%M:%S")
            line = f"[{ts}] {text}\n"
            self.log_text.configure(state="normal")
            self.log_text.insert("end", line)
            total = int(self.log_text.index("end-1c").split(".")[0])
            if total > self.MAX_LOG_LINES:
                self.log_text.delete("1.0", f"{total - self.MAX_LOG_LINES}.0")
            self.log_text.configure(state="disabled")
            if self.log_expanded:
                self._scroll_log_to_bottom()
        self.win.after(0, _do)

    def log_request(self, method, url, body=None):
        body_str = json.dumps(body, ensure_ascii=False)[:200] if body else ""
        self._append_log(f"⬆ {method} {url}", "request")
        if body_str:
            self._append_log(f"   Body: {body_str}", "request")

    def log_info(self, msg):
        self._append_log(f"ℹ️  {msg}", "info")

    def log_error(self, msg):
        self._append_log(f"❌ {msg}", "error")

    def log_success(self, msg):
        self._append_log(f"✅ {msg}", "success")

    def _scroll_log_to_bottom(self):
        self.log_text.see("end")

    # ================================================================
    #  配置页
    # ================================================================
    def _build_config_page(self):
        page = ctk.CTkScrollableFrame(self.main, fg_color="transparent")

        # API Key
        card1 = self._mk_card(page)
        card1.pack(fill="x", pady=(0, 12), padx=2)

        # 标题行：左侧标题 + 右侧「API密钥申请」按钮 + ⚠️提示
        title_row = ctk.CTkFrame(card1, fg_color="transparent")
        title_row.pack(fill="x", padx=20, pady=(18, 12))
        self._mk_label(title_row, "🔑  API 密钥", size=15, weight="bold").pack(side="left")
        # 右侧按钮组
        right_grp = ctk.CTkFrame(title_row, fg_color="transparent")
        right_grp.pack(side="right")
        btn_register = self._mk_btn(right_grp, "API密钥申请 ↗", "secondary",
                                     command=lambda: self._open_url("https://api.wlai.vip/register?aff=vtsa"))
        btn_register.pack(side="left")
        # ⚠️ 提示图标，悬浮显示申明
        warn_lbl = self._mk_label(right_grp, "⚠️", size=13, color=COLORS["warning"])
        warn_lbl.pack(side="left", padx=(6, 0))
        self._bind_tooltip(warn_lbl,
                           "申明：用该链接注册申请api密钥的同时，"
                           "云雾api会给软件开发的作者提供token的奖励，"
                           "该部分token的奖励将用于软件本身后续的升级改造，请悉知！")
        kr = ctk.CTkFrame(card1, fg_color="transparent")
        kr.pack(fill="x", padx=20, pady=(0, 6))
        self.entry_key = ctk.CTkEntry(kr, height=40, fg_color=COLORS["surface"],
                                      border_color=COLORS["divider"], border_width=1,
                                      corner_radius=8, placeholder_text="请输入云雾API的Key...",
                                      placeholder_text_color=COLORS["text_secondary"], show="•")
        self.entry_key.pack(side="left", fill="x", expand=True)
        self.entry_key.insert(0, self.config.get("api_key", ""))
        self.btn_toggle_key = self._mk_btn(kr, "👁", "ghost", command=self._toggle_key_vis)
        self.btn_toggle_key.pack(side="left", padx=(8, 0))

        self._mk_label(card1, "接口地址", size=12,
                       color=COLORS["text_secondary"]).pack(anchor="w", padx=20, pady=(10, 4))
        self.entry_base = self._mk_entry(card1)
        self.entry_base.pack(fill="x", padx=20, pady=(0, 6))
        self.entry_base.insert(0, self.config.get("api_base", ""))

        self._mk_btn(card1, "💾  保存配置", "primary",
                     command=self._save_config).pack(padx=20, pady=(8, 18), anchor="w")

        # 模型管理（查询 / 手动输入 / 多保存）
        card_model = self._mk_card(page)
        card_model.pack(fill="x", pady=(0, 12), padx=2)

        self._mk_label(card_model, "🤖  模型管理", size=15, weight="bold").pack(
            anchor="w", padx=20, pady=(18, 4))
        self._mk_label(card_model,
                       "保存你在云雾 API 内可用的模型 ID。单次/批量页的模型下拉从这里读取。\n"
                       "可点「查询可用模型」自动拉取，也可手动输入后添加。",
                       size=11, color=COLORS["text_secondary"]).pack(
            anchor="w", padx=20, pady=(0, 10))

        # 已保存模型列表
        self._mk_label(card_model, "已保存模型", size=12,
                       color=COLORS["text_secondary"]).pack(anchor="w", padx=20, pady=(0, 4))
        self.frame_saved_models = ctk.CTkFrame(card_model, fg_color=COLORS["surface"],
                                               corner_radius=8)
        self.frame_saved_models.pack(fill="x", padx=20, pady=(0, 10))

        # 查询 + 手动输入 行
        add_row = ctk.CTkFrame(card_model, fg_color="transparent")
        add_row.pack(fill="x", padx=20, pady=(0, 6))
        self.btn_query_models = self._mk_btn(add_row, "🔍 查询可用模型", "secondary",
                                             command=self._query_models)
        self.btn_query_models.pack(side="left")
        self.lbl_query_hint = self._mk_label(add_row, "", size=10,
                                             color=COLORS["text_secondary"])
        self.lbl_query_hint.pack(side="left", padx=(10, 0))

        manual_row = ctk.CTkFrame(card_model, fg_color="transparent")
        manual_row.pack(fill="x", padx=20, pady=(0, 16))
        self._mk_label(manual_row, "手动输入", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 6))
        # ⚠️ 注意符号，悬停提示模型ID须与云雾一致
        warn_lbl = self._mk_label(manual_row, "⚠️", size=13, color=COLORS["warning"])
        warn_lbl.pack(side="left", padx=(0, 6))
        self._bind_tooltip(warn_lbl,
                           "模型ID请保持与云雾api内的模型一致，否则将调用失败。")
        self.entry_manual_model = self._mk_entry(manual_row,
                                                 placeholder="例如 gpt-image-2")
        self.entry_manual_model.configure(width=200)
        self.entry_manual_model.pack(side="left", fill="x", expand=True)
        self._mk_btn(manual_row, "➕ 添加", "secondary",
                     command=self._add_manual_model).pack(side="left", padx=(8, 0))

        self._refresh_saved_models()

        # 默认路径
        card2 = self._mk_card(page)
        card2.pack(fill="x", pady=(0, 12), padx=2)

        self._mk_label(card2, "📁  默认保存路径", size=15, weight="bold").pack(
            anchor="w", padx=20, pady=(18, 12))
        pr = ctk.CTkFrame(card2, fg_color="transparent")
        pr.pack(fill="x", padx=20, pady=(0, 18))
        self.entry_save_dir = self._mk_entry(pr)
        self.entry_save_dir.pack(side="left", fill="x", expand=True)
        self.entry_save_dir.insert(0, self.config.get("save_dir", ""))

        ctk.CTkButton(pr, text="📂  浏览", height=38, width=80, corner_radius=8,
                      fg_color=COLORS["surface"], hover_color=COLORS["card_border"],
                      border_color=COLORS["divider"], border_width=1,
                      command=self._browse_save_dir).pack(side="left", padx=(8, 0))

        # 关于
        card3 = self._mk_card(page)
        card3.pack(fill="x", padx=2)

        self._mk_label(card3, "ℹ️  关于", size=15, weight="bold").pack(
            anchor="w", padx=20, pady=(18, 8))
        self._mk_label(card3,
                       "投行部智能图片生成器，基于云雾API（OpenAI 兼容协议）\n"
                       "单次出图 / 批量出图(最多30张并发) / 自定义模型管理 / gpt-image-2\n"
                       "文生图 + 图生图(上传参考图模板) / 中英文提示词 / 调用日志追踪",
                       size=12, color=COLORS["text_secondary"]).pack(
            anchor="w", padx=20, pady=(0, 10))
        self._mk_label(card3,
                       "⚠️ 注意：暂用于投行部内部使用，请勿外传",
                       size=12, weight="bold", color=COLORS["warning"]).pack(
            anchor="w", padx=20, pady=(0, 18))

        return page

    # ================================================================
    #  智能分析页（文本对话 + 思考过程 + 文件分析）
    # ================================================================
    def _build_analyze_page(self):
        page = ctk.CTkFrame(self.main, fg_color="transparent")

        # -- 参数栏（可折叠）--
        param_card = self._mk_card(page)
        param_card.pack(fill="x", pady=(0, 8), padx=2)

        # 标题栏：常用的「模型选择」常驻可见 + 折叠开关 + 清空
        prow = ctk.CTkFrame(param_card, fg_color="transparent")
        prow.pack(fill="x", padx=16, pady=(10, 8))
        self._mk_label(prow, "🤖 模型", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 6))
        self.combo_chat_model = ctk.CTkComboBox(
            prow, values=self._model_list(), width=200, height=32,
            fg_color=COLORS["surface"], button_color=COLORS["accent"],
            button_hover_color=COLORS["accent_hover"],
            border_color=COLORS["divider"], border_width=1,
            text_color=COLORS["text_primary"], font=("Microsoft YaHei UI", 12),
            command=self._on_chat_model_changed)
        self.combo_chat_model.pack(side="left", padx=(0, 4))
        # 默认选中：优先已保存的 chat_model，否则取已保存模型列表第一个
        _chat_default = self.config.get("chat_model", "")
        _models = self._model_list()
        if _chat_default not in _models:
            _chat_default = _models[0] if _models else "gpt-4o"
        self.combo_chat_model.set(_chat_default)
        self.lbl_chat_reason = self._mk_label(prow, "", size=10,
                                              color=COLORS["accent"])
        self.lbl_chat_reason.pack(side="left", padx=(2, 0))

        self._mk_btn(prow, "🗑 清空", "ghost",
                     command=self._clear_chat).pack(side="right")
        # 折叠/展开「高级参数」的按钮（默认收起，给对话区让出空间）
        self.chat_params_open = False
        self.btn_toggle_params = ctk.CTkButton(
            prow, text="⚙ 参数 ▾", height=30, width=92, corner_radius=8,
            fg_color="transparent", hover_color=COLORS["surface"],
            border_color=COLORS["divider"], border_width=1,
            text_color=COLORS["text_secondary"], font=("Microsoft YaHei UI", 11),
            command=self._toggle_chat_params)
        self.btn_toggle_params.pack(side="right", padx=(0, 8))

        # 折叠容器：温度/最大回复/流式 + System + 文件上传，全部放这里
        self.chat_param_body = ctk.CTkFrame(param_card, fg_color="transparent")
        # 默认不 pack（收起）

        # 第二行：温度 / 最大回复 / 流式
        prow2 = ctk.CTkFrame(self.chat_param_body, fg_color="transparent")
        prow2.pack(fill="x", padx=16, pady=(0, 6))
        self._mk_label(prow2, "温度", size=11,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 2))
        self._info_dot(
            prow2,
            "温度 (temperature)：控制回复的随机性。\n"
            "· 取值 0–2，默认 0.7\n"
            "· 越低越严谨、确定、可复现，适合代码/事实问答\n"
            "· 越高越发散、有创意，适合头脑风暴/文案创作").pack(
            side="left", padx=(0, 4))
        self.lbl_temp_val = self._mk_label(prow2, "", size=11,
                                           color=COLORS["accent"], width=30)
        self.temp_slider = ctk.CTkSlider(
            prow2, from_=0, to=2, number_of_steps=20, width=130,
            fg_color=COLORS["surface"], progress_color=COLORS["accent"],
            button_color=COLORS["accent"], button_hover_color=COLORS["accent_hover"],
            command=self._on_temp_changed)
        self.temp_slider.pack(side="left", padx=(0, 2))
        self.temp_slider.set(float(self.config.get("chat_temperature", 0.7)))
        self.lbl_temp_val.pack(side="left", padx=(0, 12))
        self._on_temp_changed(self.temp_slider.get())

        self._mk_label(prow2, "最大回复", size=11,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 2))
        self._info_dot(
            prow2,
            "最大回复 (max_tokens)：单次回复的最大长度上限。\n"
            "· 单位为 token，约 1 个汉字≈1.5 token\n"
            "· 设太小会导致长回复被中途截断\n"
            "· 设太大不会浪费，模型按需生成；建议 2048 起步").pack(
            side="left", padx=(0, 4))
        self.entry_chat_maxtok = ctk.CTkEntry(
            prow2, width=70, height=30, fg_color=COLORS["surface"],
            border_color=COLORS["divider"], border_width=1, corner_radius=6,
            text_color=COLORS["text_primary"], font=("Consolas", 11))
        self.entry_chat_maxtok.pack(side="left", padx=(0, 12))
        self.entry_chat_maxtok.insert(0, str(self.config.get("chat_max_tokens", 2048)))

        self.chat_stream_var = ctk.BooleanVar(value=self.config.get("chat_stream", True))
        ctk.CTkCheckBox(
            prow2, text=" 流式输出", variable=self.chat_stream_var,
            checkbox_width=16, checkbox_height=16, corner_radius=4,
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            border_color=COLORS["divider"], border_width=2,
            text_color=COLORS["text_primary"],
            font=("Microsoft YaHei UI", 11)).pack(side="left", padx=(0, 2))
        self._info_dot(
            prow2,
            "流式输出 (stream)：是否边生成边逐字显示。\n"
            "· 开启：像打字机一样实时显示，等待感更低（推荐）\n"
            "· 关闭：等模型完全生成后一次性显示\n"
            "· 推理模型的「思考过程」仅在流式下可实时查看").pack(
            side="left")

        # System 提示词
        srow = ctk.CTkFrame(self.chat_param_body, fg_color="transparent")
        srow.pack(fill="x", padx=16, pady=(0, 8))
        self._mk_label(srow, "System", size=11,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 6))
        self.entry_chat_system = ctk.CTkEntry(
            srow, height=30, fg_color=COLORS["surface"],
            border_color=COLORS["divider"], border_width=1, corner_radius=6,
            text_color=COLORS["text_primary"], font=("Microsoft YaHei UI", 11),
            placeholder_text="系统提示词（设定助手角色）",
            placeholder_text_color=COLORS["text_secondary"])
        self.entry_chat_system.pack(side="left", fill="x", expand=True)
        self.entry_chat_system.insert(0, self.config.get("chat_system", ""))

        # 文件分析行：上传任意格式文件 + 已选文件芯片
        frow = ctk.CTkFrame(self.chat_param_body, fg_color="transparent")
        frow.pack(fill="x", padx=16, pady=(0, 12))
        ctk.CTkButton(
            frow, text="📎 上传文件分析", height=30, width=130, corner_radius=8,
            fg_color=COLORS["surface"], hover_color=COLORS["card_border"],
            border_color=COLORS["divider"], border_width=1,
            text_color=COLORS["text_primary"], font=("Microsoft YaHei UI", 11),
            command=self._add_chat_files).pack(side="left", padx=(0, 8))
        self.chat_files_frame = ctk.CTkFrame(frow, fg_color="transparent")
        self.chat_files_frame.pack(side="left", fill="x", expand=True)
        self._mk_label(
            frow, "支持任意格式（文本类直接解析内容）", size=10,
            color=COLORS["text_secondary"]).pack(side="right")

        # -- 对话消息区（滚动）--
        self.chat_scroll = ctk.CTkScrollableFrame(
            page, fg_color=COLORS["card"], corner_radius=12)
        self.chat_scroll.pack(fill="both", expand=True, pady=(0, 8), padx=2)
        self.chat_placeholder = self._mk_label(
            self.chat_scroll,
            "💬 在下方输入消息，开始对话\n"
            "推理模型(o1/o3/deepseek-reasoner 等)会展示思考过程\n"
            "可上传任意格式文件，文本内容会自动随消息一起发送\n"
            "每条回复右下角可「复制」或「查看全文」",
            size=13, color=COLORS["text_secondary"])
        self.chat_placeholder.pack(pady=40)

        # -- 输入区 --
        input_card = self._mk_card(page)
        input_card.pack(fill="x", padx=2)
        irow = ctk.CTkFrame(input_card, fg_color="transparent")
        irow.pack(fill="x", padx=12, pady=12)
        self.text_chat_input = ctk.CTkTextbox(
            irow, height=64, fg_color=COLORS["surface"],
            border_color=COLORS["divider"], border_width=1, corner_radius=8,
            wrap="word", font=("Microsoft YaHei UI", 13),
            text_color=COLORS["text_primary"])
        self.text_chat_input.pack(side="left", fill="both", expand=True, padx=(0, 8))
        self.text_chat_input.bind("<Return>", self._on_chat_enter)
        self.text_chat_input.bind("<Shift-Return>", lambda e: None)

        self.btn_chat_send = ctk.CTkButton(
            irow, text="发送", width=72, height=64, corner_radius=10,
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            font=("Microsoft YaHei UI", 14, "bold"),
            command=self._send_chat)
        self.btn_chat_send.pack(side="left")

        self._on_chat_model_changed(self.combo_chat_model.get())
        return page

    def _toggle_chat_params(self):
        """展开/收起「高级参数」区域，收起时把空间让给对话区。"""
        self.chat_params_open = not self.chat_params_open
        if self.chat_params_open:
            # chat_param_body 的父级是 param_card，直接 pack 会排在标题行 prow 之后
            self.chat_param_body.pack(fill="x")
            self.btn_toggle_params.configure(text="⚙ 参数 ▴")
        else:
            self.chat_param_body.pack_forget()
            self.btn_toggle_params.configure(text="⚙ 参数 ▾")

    def _on_temp_changed(self, val):
        self.lbl_temp_val.configure(text=f"{float(val):.1f}")

    def _on_chat_model_changed(self, choice):
        self.config["chat_model"] = choice
        if is_reasoning_model(choice):
            self.lbl_chat_reason.configure(text="🧠 推理模型（显示思考过程）")
        else:
            self.lbl_chat_reason.configure(text="")

    def _on_chat_enter(self, event):
        if event.state & 0x0001:   # Shift+回车换行
            return
        self._send_chat()
        return "break"

    # ---- 文件分析 ----
    def _add_chat_files(self):
        paths = filedialog.askopenfilenames(title="选择要分析的文件（任意格式）")
        if not paths:
            return
        for p in paths:
            if any(f["path"] == p for f in self.chat_files):
                continue
            ext = os.path.splitext(p)[1].lower()
            self.chat_files.append({"path": p, "name": os.path.basename(p), "ext": ext})
        self._render_chat_file_chips()
        self.log_info(f"已添加待分析文件，共 {len(self.chat_files)} 个")

    def _render_chat_file_chips(self):
        for w in self.chat_files_frame.winfo_children():
            w.destroy()
        for f in self.chat_files:
            chip = ctk.CTkFrame(self.chat_files_frame, fg_color=COLORS["surface"],
                                corner_radius=6)
            chip.pack(side="left", padx=(0, 6))
            readable = f["ext"] in TEXT_READABLE_EXTS or f["ext"] in BINARY_DOC_EXTS
            tag = "📄" if readable else "📦"
            name = f["name"]
            if len(name) > 18:
                name = name[:15] + "…"
            self._mk_label(chip, f"{tag} {name}", size=10,
                           color=COLORS["text_primary"]).pack(side="left", padx=(8, 2),
                                                              pady=4)
            ctk.CTkButton(
                chip, text="✕", width=18, height=18, corner_radius=4,
                fg_color="transparent", hover_color=COLORS["card_border"],
                text_color=COLORS["text_secondary"], font=("Microsoft YaHei UI", 10),
                command=lambda fp=f["path"]: self._remove_chat_file(fp)).pack(
                side="left", padx=(0, 4))

    def _remove_chat_file(self, path):
        self.chat_files = [f for f in self.chat_files if f["path"] != path]
        self._render_chat_file_chips()

    def _build_files_context(self):
        """把已上传文件读成文本上下文。文本类直接读内容，二进制文档用专用库解析。
        返回拼接好的上下文字符串（可能为空）。"""
        if not self.chat_files:
            return ""
        parts = []
        for f in self.chat_files:
            p, name, ext = f["path"], f["name"], f["ext"]
            if ext in TEXT_READABLE_EXTS:
                # ---- 纯文本文件：直接读取 ----
                try:
                    with open(p, "r", encoding="utf-8", errors="replace") as fh:
                        content = fh.read(FILE_READ_MAX)
                    if os.path.getsize(p) > FILE_READ_MAX:
                        content += "\n…(内容过长，仅截取前 200KB)…"
                    parts.append(f"【文件：{name}】\n{content}")
                except Exception as e:
                    parts.append(f"【文件：{name}】读取失败：{e}")
            elif ext == ".docx":
                # ---- Word 文档：python-docx ----
                try:
                    from docx import Document
                    doc = Document(p)
                    paragraphs = [pp.text for pp in doc.paragraphs if pp.text.strip()]
                    for table in doc.tables:
                        rows = []
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            rows.append(" | ".join(cells))
                        if rows:
                            paragraphs.append("（表格）\n" + "\n".join(rows))
                    content = "\n".join(paragraphs)
                    if len(content.encode('utf-8')) > FILE_READ_MAX:
                        content = content[:FILE_READ_MAX] + "\n…(内容过长，仅截取前 200KB)…"
                    parts.append(f"【文件：{name}（Word 文档）】\n{content}")
                except Exception as e:
                    parts.append(f"【文件：{name}】读取失败：{e}")
            elif ext == ".xlsx":
                # ---- Excel 表格：openpyxl ----
                try:
                    from openpyxl import load_workbook
                    wb = load_workbook(p, read_only=True, data_only=True)
                    all_text = []
                    for sheet_name in wb.sheetnames:
                        ws = wb[sheet_name]
                        rows_text = []
                        for row in ws.iter_rows(values_only=True):
                            cells = [str(c) if c is not None else "" for c in row]
                            if any(cells):  # 跳过全空行
                                rows_text.append(" | ".join(cells))
                        if rows_text:
                            all_text.append(f"--- 工作表：{sheet_name} ---\n" + "\n".join(rows_text))
                    wb.close()
                    content = "\n\n".join(all_text)
                    if len(content.encode('utf-8')) > FILE_READ_MAX:
                        content = content[:FILE_READ_MAX] + "\n…(内容过长，仅截取前 200KB)…"
                    parts.append(f"【文件：{name}（Excel 表格）】\n{content}")
                except Exception as e:
                    parts.append(f"【文件：{name}】读取失败：{e}")
            elif ext == ".pptx":
                # ---- PPT 幻灯片：python-pptx ----
                try:
                    from pptx import Presentation
                    prs = Presentation(p)
                    all_text = []
                    for i, slide in enumerate(prs.slides, 1):
                        slide_text = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    text = para.text.strip()
                                    if text:
                                        slide_text.append(text)
                            if shape.has_table:
                                for row in shape.table.rows:
                                    cells = [cell.text.strip() for cell in row.cells]
                                    slide_text.append(" | ".join(cells))
                        if slide_text:
                            all_text.append(f"--- 第 {i} 页 ---\n" + "\n".join(slide_text))
                    content = "\n\n".join(all_text)
                    if len(content.encode('utf-8')) > FILE_READ_MAX:
                        content = content[:FILE_READ_MAX] + "\n…(内容过长，仅截取前 200KB)…"
                    parts.append(f"【文件：{name}（PPT 幻灯片）】\n{content}")
                except Exception as e:
                    parts.append(f"【文件：{name}】读取失败：{e}")
            elif ext == ".pdf":
                # ---- PDF 文档：PyPDF2 ----
                try:
                    from PyPDF2 import PdfReader
                    reader = PdfReader(p)
                    all_text = []
                    total_pages = len(reader.pages)
                    for i, page in enumerate(reader.pages):
                        text = page.extract_text()
                        if text and text.strip():
                            all_text.append(f"--- 第 {i+1}/{total_pages} 页 ---\n{text.strip()}")
                    content = "\n\n".join(all_text)
                    if len(content.encode('utf-8')) > FILE_READ_MAX:
                        content = content[:FILE_READ_MAX] + "\n…(内容过长，仅截取前 200KB)…"
                    parts.append(f"【文件：{name}（PDF 文档，共 {total_pages} 页）】\n{content}")
                except Exception as e:
                    parts.append(f"【文件：{name}】读取失败：{e}")
            else:
                # ---- 其他二进制格式：只附元信息 ----
                try:
                    size_kb = os.path.getsize(p) / 1024
                except Exception:
                    size_kb = 0
                parts.append(
                    f"【文件：{name}】（{ext or '未知'} 格式，{size_kb:.0f}KB，"
                    f"非文本格式无法直接解析内容，请根据文件名/类型进行说明或建议）")
        return "以下是用户上传的待分析文件内容：\n\n" + "\n\n".join(parts)

    def _adjust_textbox_height(self, text_box, text):
        """根据文本内容自适应调整 CTkTextbox 的高度。
        按换行符和 wrap 宽度估算行数，每行约 20px 高度。"""
        if not text:
            text_box.configure(height=30)
            return
        # 估算可用宽度（像素）：气泡宽度减去左右 padding
        # bubble 最大宽度约 700px，text_box padx=8 两边共 16px
        avail_width = 660
        # 中文字符约 13px 宽，英文约 7px 宽，保守按平均每字符 10px
        chars_per_line = max(20, avail_width // 10)
        lines = text.split('\n')
        total_lines = 0
        for line in lines:
            line_len = len(line)
            # 空行也算一行
            wrapped = max(1, (line_len + chars_per_line - 1) // chars_per_line)
            total_lines += wrapped
        # 每行高度约 20px，最小 30px，最大 800px（避免过长）
        height = max(30, min(800, total_lines * 20 + 10))
        text_box.configure(height=height)

    # ---- 对话气泡 ----
    def _clear_chat(self):
        if self.chat_streaming:
            return
        self.chat_history = []
        for w in self.chat_scroll.winfo_children():
            w.destroy()
        self.chat_placeholder = self._mk_label(
            self.chat_scroll, "💬 已清空，开始新的对话",
            size=13, color=COLORS["text_secondary"])
        self.chat_placeholder.pack(pady=40)

    def _add_chat_bubble(self, role, text=""):
        if self.chat_placeholder is not None:
            self.chat_placeholder.destroy()
            self.chat_placeholder = None
        is_user = (role == "user")
        wrapper = ctk.CTkFrame(self.chat_scroll, fg_color="transparent")
        wrapper.pack(fill="x", padx=10, pady=(6, 2))
        bubble = ctk.CTkFrame(
            wrapper, fg_color=(COLORS["accent"] if is_user else COLORS["surface"]),
            corner_radius=12)
        # 助手气泡占更宽（右留 40），用户气泡靠右（左留 80）
        bubble.pack(anchor="e" if is_user else "w",
                    fill="x" if not is_user else None,
                    padx=(80, 0) if is_user else (0, 40))
        prefix = "🧑 你" if is_user else "🤖 助手"
        self._mk_label(bubble, prefix, size=10,
                       color="white" if is_user else COLORS["text_secondary"]).pack(
            anchor="w", padx=12, pady=(8, 0))

        # 使用 CTkTextbox 替代 CTkLabel，更好地保留段落和换行
        if is_user:
            # 用户消息仍用 Label（通常较短）
            lbl = ctk.CTkLabel(
                bubble, text=text, font=("Microsoft YaHei UI", 13),
                text_color="white",
                justify="left", wraplength=520)
            lbl.pack(anchor="w", padx=12, pady=(2, 10))
        else:
            # 助手回复用 Textbox，保留段落格式
            text_box = ctk.CTkTextbox(
                bubble,
                fg_color="transparent",
                border_width=0,
                corner_radius=0,
                wrap="word",
                font=("Microsoft YaHei UI", 13),
                text_color=COLORS["text_primary"],
                activate_scrollbars=False,
            )
            text_box.pack(anchor="w", fill="x", padx=8, pady=(2, 6))
            text_box.insert("1.0", text)
            text_box.configure(state="disabled")
            # 自适应高度：根据文本行数计算
            self._adjust_textbox_height(text_box, text)
            # 绑定标签用于复制/查看全文
            text_box._text = text
            lbl = text_box  # 保持返回类型一致

        # 助手气泡：底部操作栏（复制 / 查看全文）
        if not is_user:
            actions = ctk.CTkFrame(bubble, fg_color="transparent")
            actions.pack(anchor="e", fill="x", padx=10, pady=(0, 8))
            ctk.CTkButton(
                actions, text="📋 复制", height=24, width=64, corner_radius=6,
                fg_color="transparent", hover_color=COLORS["card_border"],
                border_color=COLORS["divider"], border_width=1,
                text_color=COLORS["text_secondary"], font=("Microsoft YaHei UI", 10),
                command=lambda l=lbl: self._copy_chat_text(l)).pack(side="right")
            ctk.CTkButton(
                actions, text="⤢ 查看全文", height=24, width=84, corner_radius=6,
                fg_color="transparent", hover_color=COLORS["card_border"],
                border_color=COLORS["divider"], border_width=1,
                text_color=COLORS["text_secondary"], font=("Microsoft YaHei UI", 10),
                command=lambda l=lbl: self._view_chat_fulltext(l)).pack(
                side="right", padx=(0, 6))

        self._chat_scroll_bottom()
        return lbl

    def _copy_chat_text(self, lbl):
        """复制某条回复的完整文本到剪贴板。"""
        # 兼容 CTkLabel 和 CTkTextbox
        if hasattr(lbl, '_text'):
            text = lbl._text or ""
        else:
            text = lbl.cget("text") or ""
        try:
            self.win.clipboard_clear()
            self.win.clipboard_append(text)
            self._set_status(f"✅ 已复制 {len(text)} 字到剪贴板", COLORS["success"])
        except Exception as e:
            self._set_status(f"复制失败：{e}", COLORS["error"])

    def _view_chat_fulltext(self, lbl):
        """弹出独立窗口查看回复全文，可自由滚动、全选复制，解决小框滚动的不适感。"""
        # 兼容 CTkLabel 和 CTkTextbox
        if hasattr(lbl, '_text'):
            text = lbl._text or ""
        else:
            text = lbl.cget("text") or ""
        win = ctk.CTkToplevel(self.win)
        win.title("回复全文")
        win.geometry("720x600")
        win.configure(fg_color=COLORS["bg"])
        win.transient(self.win)

        top = ctk.CTkFrame(win, fg_color="transparent")
        top.pack(fill="x", padx=16, pady=(14, 6))
        self._mk_label(top, f"回复全文（{len(text)} 字）", size=14,
                       weight="bold").pack(side="left")
        ctk.CTkButton(
            top, text="📋 复制全部", height=30, width=96, corner_radius=8,
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            font=("Microsoft YaHei UI", 11),
            command=lambda: self._copy_chat_text(lbl)).pack(side="right")

        box = ctk.CTkTextbox(
            win, fg_color=COLORS["surface"], border_width=1,
            border_color=COLORS["divider"], corner_radius=10, wrap="word",
            font=("Microsoft YaHei UI", 13), text_color=COLORS["text_primary"])
        box.pack(fill="both", expand=True, padx=16, pady=(0, 16))
        box.insert("1.0", text)
        box.configure(state="disabled")
        win.after(50, win.lift)

    def _add_reason_box(self):
        """推理模型：添加折叠的「思考过程」区域，返回 textbox。"""
        wrapper = ctk.CTkFrame(self.chat_scroll, fg_color="transparent")
        wrapper.pack(fill="x", padx=10, pady=(6, 0))
        box_frame = ctk.CTkFrame(wrapper, fg_color=COLORS["bg"], corner_radius=10,
                                 border_width=1, border_color=COLORS["divider"])
        box_frame.pack(anchor="w", fill="x", padx=(0, 60))
        self._mk_label(box_frame, "🧠 思考过程", size=10,
                       color=COLORS["text_secondary"]).pack(anchor="w", padx=12,
                                                            pady=(8, 2))
        box = ctk.CTkTextbox(
            box_frame, height=90, fg_color=COLORS["bg"], border_width=0,
            wrap="word", font=("Microsoft YaHei UI", 11),
            text_color=COLORS["text_secondary"])
        box.pack(fill="x", padx=8, pady=(0, 8))
        box.configure(state="disabled")
        return box

    def _chat_scroll_bottom(self):
        """节流滚动到底部：100ms 内合并多次调用，避免流式追加时频繁滚动卡顿。"""
        if hasattr(self, '_scroll_bottom_pending') and self._scroll_bottom_pending:
            return  # 已有待执行的滚动，跳过
        self._scroll_bottom_pending = True
        def _do_scroll():
            self._scroll_bottom_pending = False
            try:
                self.chat_scroll._parent_canvas.yview_moveto(1.0)
            except Exception:
                pass
        self.win.after(100, _do_scroll)

    def _send_chat(self):
        if self.chat_streaming:
            return
        msg = self.text_chat_input.get("1.0", "end-1c").strip()
        if not msg and not self.chat_files:
            return
        oa_base = self.config.get("api_base", "").strip()
        oa_key = self.config.get("api_key", "").strip()
        if not oa_base or not oa_key:
            messagebox.showwarning(
                "提示", "文本对话需要 OpenAI 协议接口，请到「配置」页填写接口地址与 API Key")
            self._switch_tab("config")
            return

        model = self.combo_chat_model.get().strip()
        try:
            temperature = float(self.temp_slider.get())
        except Exception:
            temperature = 0.7
        try:
            max_tokens = int(self.entry_chat_maxtok.get().strip())
        except Exception:
            max_tokens = 2048
        system = self.entry_chat_system.get().strip()
        stream = bool(self.chat_stream_var.get())

        self.config.update({
            "chat_model": model, "chat_system": system,
            "chat_temperature": temperature, "chat_max_tokens": max_tokens,
            "chat_stream": stream,
        })
        save_config(self.config)

        # 组装本轮用户消息：文件上下文 + 文本问题
        file_ctx = self._build_files_context()
        if file_ctx:
            user_for_api = file_ctx + "\n\n【我的问题】\n" + (msg or "请分析以上文件内容。")
            bubble_text = (msg or "请分析以上文件内容。") + \
                f"\n\n📎 已附 {len(self.chat_files)} 个文件"
        else:
            user_for_api = msg
            bubble_text = msg

        self._add_chat_bubble("user", bubble_text)
        self.chat_history.append({"role": "user", "content": user_for_api})
        self.text_chat_input.delete("1.0", "end")
        # 文件随本轮提交后清空（已并入历史）
        self.chat_files = []
        self._render_chat_file_chips()

        messages = []
        if system:
            messages.append({"role": "system", "content": system})
        messages.extend(self.chat_history)

        reason_box = self._add_reason_box() if is_reasoning_model(model) else None
        content_lbl = self._add_chat_bubble("assistant", "")
        self._chat_cur_content_lbl = content_lbl
        self._chat_cur_reason_box = reason_box
        self._chat_cur_content_buf = ""
        self._chat_cur_reason_buf = ""
        # 追加模式标志：流式输出时用 insert("end", delta) 而非全删全插
        self._chat_content_appended = 0   # 已追加到 Textbox 的字符数
        self._chat_reason_appended = 0     # 已追加到 reason_box 的字符数

        self.chat_streaming = True
        self.chat_stop = False
        self.btn_chat_send.configure(text="生成中", state="disabled",
                                     fg_color=COLORS["surface"],
                                     text_color=COLORS["text_secondary"])
        t = threading.Thread(
            target=self._chat_worker,
            args=(oa_base, oa_key, model, messages, temperature, max_tokens, stream),
            daemon=True)
        t.start()

    def _chat_worker(self, oa_base, oa_key, model, messages,
                     temperature, max_tokens, stream):
        try:
            self.log_request("POST",
                             f"{_normalize_openai_base(oa_base)}/chat/completions",
                             {"model": model, "temperature": temperature,
                              "max_tokens": max_tokens, "stream": stream})

            def on_reason(d):
                self._chat_cur_reason_buf += d
                if stream:
                    # 流式追加模式：只插入新增的 delta，避免全删全插的闪烁
                    self._ui(lambda delta=d: self._append_reason_box(delta))
                else:
                    self._ui(lambda: self._update_reason_box(self._chat_cur_reason_buf))

            def on_content(d):
                self._chat_cur_content_buf += d
                if stream:
                    # 流式追加模式：只插入新增的 delta
                    self._ui(lambda delta=d: self._append_content_lbl(delta))
                else:
                    self._ui(lambda: self._update_content_lbl(self._chat_cur_content_buf))

            if stream:
                reason, content = openai_chat_stream(
                    oa_base, oa_key, model, messages,
                    temperature=temperature, max_tokens=max_tokens,
                    on_reasoning=on_reason, on_content=on_content,
                    stop_flag=lambda: self.chat_stop)
            else:
                reason, content = openai_chat_once(
                    oa_base, oa_key, model, messages,
                    temperature=temperature, max_tokens=max_tokens)
                if reason:
                    self._ui(lambda r=reason: self._update_reason_box(r))
                self._ui(lambda c=content: self._update_content_lbl(c))

            final = content or "(空响应)"
            self.chat_history.append({"role": "assistant", "content": final})
            self.log_success(f"回复完成 | {len(final)} 字"
                             + (f" | 思考 {len(reason)} 字" if reason else ""))
        except Exception as e:
            err = str(e)
            self._ui(lambda: self._update_content_lbl(f"❌ {err[:200]}"))
            self.log_error(f"回复失败: {err[:120]}")
        finally:
            self._ui(self._chat_done)

    def _update_content_lbl(self, text):
        """全量替换文本（非流式 / 错误信息 / 最终刷新时使用）。"""
        if self._chat_cur_content_lbl is not None:
            lbl = self._chat_cur_content_lbl
            # 兼容 CTkLabel 和 CTkTextbox
            if hasattr(lbl, '_text'):
                # CTkTextbox
                lbl.configure(state="normal")
                lbl.delete("1.0", "end")
                lbl.insert("1.0", text)
                lbl.configure(state="disabled")
                lbl._text = text
                # 自适应调整高度
                self._adjust_textbox_height(lbl, text)
            else:
                # CTkLabel
                lbl.configure(text=text)
            self._chat_scroll_bottom()

    def _append_content_lbl(self, delta):
        """流式追加模式：只插入新增的 delta 文本，避免全删全插的闪烁。"""
        if self._chat_cur_content_lbl is not None:
            lbl = self._chat_cur_content_lbl
            if hasattr(lbl, '_text'):
                # CTkTextbox — 追加 delta
                lbl.configure(state="normal")
                lbl.insert("end", delta)
                lbl.configure(state="disabled")
                lbl._text = (lbl._text or "") + delta
                self._chat_content_appended += len(delta)
                # 每 ~200 字符调整一次高度（避免频繁重排卡顿）
                if self._chat_content_appended % 200 < len(delta):
                    self._adjust_textbox_height(lbl, lbl._text)
            else:
                # CTkLabel — 退回全量替换
                lbl.configure(text=(lbl.cget("text") or "") + delta)
            self._chat_scroll_bottom()
            # 节流刷新 UI：50ms 内合并多次 update_idletasks
            if not hasattr(self, '_ui_flush_pending') or not self._ui_flush_pending:
                self._ui_flush_pending = True
                def _flush():
                    self._ui_flush_pending = False
                    try:
                        self.win.update_idletasks()
                    except Exception:
                        pass
                self.win.after(50, _flush)

    def _update_reason_box(self, text):
        """全量替换推理文本（非流式时使用）。"""
        box = self._chat_cur_reason_box
        if box is not None:
            box.configure(state="normal")
            box.delete("1.0", "end")
            box.insert("1.0", text)
            box.configure(state="disabled")
            # 推理文本较长时自适应高度
            self._adjust_textbox_height(box, text)
            self._chat_scroll_bottom()

    def _append_reason_box(self, delta):
        """流式追加推理文本 delta，避免全删全插。"""
        box = self._chat_cur_reason_box
        if box is not None:
            box.configure(state="normal")
            box.insert("end", delta)
            box.configure(state="disabled")
            self._chat_reason_appended += len(delta)
            # 每 ~200 字符调整高度
            if self._chat_reason_appended % 200 < len(delta):
                full = self._chat_cur_reason_buf
                self._adjust_textbox_height(box, full)
            self._chat_scroll_bottom()

    def _chat_done(self):
        self.chat_streaming = False
        self.btn_chat_send.configure(text="发送", state="normal",
                                     fg_color=COLORS["accent"], text_color="white")
        # 最终全量刷新：确保高度计算精确 + _text 缓存与显示一致
        if self._chat_cur_content_lbl is not None and hasattr(self._chat_cur_content_lbl, '_text'):
            self._adjust_textbox_height(self._chat_cur_content_lbl,
                                        self._chat_cur_content_lbl._text or "")
        if self._chat_cur_reason_box is not None and self._chat_cur_reason_buf:
            self._adjust_textbox_height(self._chat_cur_reason_box,
                                        self._chat_cur_reason_buf)
        self._chat_cur_content_lbl = None
        self._chat_cur_reason_box = None

    # ================================================================
    #  单次出图页
    # ================================================================
    def _build_single_page(self):
        page = ctk.CTkScrollableFrame(self.main, fg_color="transparent")

        # -- 参考图附件（紧凑卡片，单次/批量共用同一构建器）--
        self._build_ref_card(page, self.ctx_single)

        # -- 提示词（单条） --
        card_prompt = self._mk_card(page)
        card_prompt.pack(fill="x", pady=(0, 10), padx=2)

        self._mk_label(card_prompt, "✏️  提示词 Prompt", size=14, weight="bold").pack(
            anchor="w", padx=20, pady=(16, 4))
        self.text_single_prompt = ctk.CTkTextbox(
            card_prompt, height=150, fg_color=COLORS["surface"],
            border_color=COLORS["divider"], border_width=1, corner_radius=8,
            wrap="word", font=("Microsoft YaHei UI", 14),
            text_color=COLORS["text_primary"],
        )
        self.text_single_prompt.pack(fill="x", padx=20, pady=(4, 16))

        # -- 参数设置 --
        card_params = self._mk_card(page)
        card_params.pack(fill="x", pady=(0, 10), padx=2)

        self._mk_label(card_params, "⚙️  参数设置", size=14, weight="bold").pack(
            anchor="w", padx=20, pady=(16, 10))

        # 模型
        rm = ctk.CTkFrame(card_params, fg_color="transparent")
        rm.pack(fill="x", padx=20, pady=(0, 6))
        self._mk_label(rm, "模型", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 10))
        self.combo_single_model = self._mk_optionmenu(
            rm, self._model_list(), width=240,
            command=self._on_single_model_changed)
        self.combo_single_model.pack(side="left")
        self.combo_single_model.set(self._current_model())
        self.lbl_single_model_hint = self._mk_label(rm, "", size=10,
                                                    color=COLORS["text_secondary"])
        self.lbl_single_model_hint.pack(side="left", padx=(10, 0))

        # 尺寸
        r1 = ctk.CTkFrame(card_params, fg_color="transparent")
        r1.pack(fill="x", padx=20, pady=(0, 6))
        self._mk_label(r1, "尺寸比例", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 10))
        self.combo_single_size = self._mk_optionmenu(
            r1, list(SIZE_MAP.keys()), width=170,
            command=self._on_single_size_changed)
        self.combo_single_size.pack(side="left")
        lk = self.config.get("last_size_key", DEFAULT_SIZE_KEY)
        self.combo_single_size.set(lk if lk in SIZE_MAP else DEFAULT_SIZE_KEY)
        self.lbl_single_pixel = self._mk_label(
            r1, self._size_caption(self.combo_single_size.get(), self.combo_single_model.get()),
            size=11, color=COLORS["text_secondary"])
        self.lbl_single_pixel.pack(side="left", padx=(10, 0))

        # 文件名
        r2 = ctk.CTkFrame(card_params, fg_color="transparent")
        r2.pack(fill="x", padx=20, pady=(0, 6))
        self._mk_label(r2, "文件名", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 10))
        self.entry_single_prefix = self._mk_entry(r2)
        self.entry_single_prefix.configure(width=160)
        self.entry_single_prefix.pack(side="left")
        self.entry_single_prefix.insert(0, self.config.get("filename_prefix", "ai_image"))
        self._mk_label(r2, ".png（自动追加时间戳）", size=10,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(8, 0))

        # 出图画质（gpt-image-2 quality 档位，单次页可临时切换）
        rq = ctk.CTkFrame(card_params, fg_color="transparent")
        rq.pack(fill="x", padx=20, pady=(0, 6))
        self._mk_label(rq, "出图画质", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 10))
        self.combo_single_quality = self._mk_optionmenu(
            rq, QUALITY_OPTIONS, width=120,
            command=self._on_single_quality_changed)
        self.combo_single_quality.pack(side="left")
        self.combo_single_quality.set(self._quality())
        self._mk_label(rq, "（auto 由模型自动决定；high 更慢更贵）", size=10,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(8, 0))

        # 保存路径提示
        r3 = ctk.CTkFrame(card_params, fg_color="transparent")
        r3.pack(fill="x", padx=20, pady=(0, 14))
        self._mk_label(r3, "保存到", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 10))
        self.lbl_single_out_dir = self._mk_label(
            r3, self.config.get("save_dir", "") or "（未设置，请到「配置」页设置）",
            size=12, color=COLORS["text_primary"])
        self.lbl_single_out_dir.pack(side="left")

        # -- 生成按钮 --
        bf = ctk.CTkFrame(page, fg_color="transparent")
        bf.pack(fill="x", padx=2, pady=(6, 0))
        self.ctx_single.btn_generate = ctk.CTkButton(
            bf, text="🚀  开始生成", height=48, corner_radius=12,
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            font=("Microsoft YaHei UI", 15, "bold"),
            command=self._start_single)
        self.ctx_single.btn_generate.pack(fill="x")

        # -- 进度 + 一览 --
        self._build_progress_and_gallery(page, self.ctx_single,
                                         "✨ 图片将直接生成并保存到默认路径，生成后在此查看")

        self._refresh_single_model_hint()
        return page

    # ================================================================
    #  批量出图页
    # ================================================================
    def _build_batch_page(self):
        page = ctk.CTkScrollableFrame(self.main, fg_color="transparent")

        # -- 参考图附件（置顶，与单次出图风格统一；作用于本批所有提示词）--
        self._build_ref_card(page, self.ctx_batch,
                             note="（作为模板应用到本批所有提示词）",
                             style_lock=True)

        # -- 提示词列表 --
        card_list = self._mk_card(page)
        card_list.pack(fill="x", pady=(0, 10), padx=2)

        lh = ctk.CTkFrame(card_list, fg_color="transparent")
        lh.pack(fill="x", padx=20, pady=(16, 4))
        self._mk_label(lh, "📝  提示词列表（每条单独生成 1 张图）", size=14,
                       weight="bold").pack(side="left")
        self._mk_btn(lh, "📥  导入txt", "secondary",
                     command=self._import_batch_prompts).pack(side="right")
        self._mk_btn(lh, "🗑  清空", "ghost",
                     command=self._clear_batch_rows).pack(side="right", padx=(0, 8))

        # 行容器
        self.batch_rows_frame = ctk.CTkFrame(card_list, fg_color="transparent")
        self.batch_rows_frame.pack(fill="x", padx=20, pady=(4, 4))

        # 底部：添加一行 + 计数
        add_row = ctk.CTkFrame(card_list, fg_color="transparent")
        add_row.pack(fill="x", padx=20, pady=(4, 16))
        self._mk_btn(add_row, "➕  添加一行", "secondary",
                     command=lambda: self._add_batch_row()).pack(side="left")
        self.lbl_batch_count = self._mk_label(add_row, "", size=12,
                                              color=COLORS["accent"])
        self.lbl_batch_count.pack(side="left", padx=(16, 0))

        # 默认 5 行
        for _ in range(BATCH_DEFAULT_ROWS):
            self._add_batch_row()

        # -- 参数设置 --
        card_params = self._mk_card(page)
        card_params.pack(fill="x", pady=(0, 10), padx=2)

        self._mk_label(card_params, "⚙️  参数设置", size=14, weight="bold").pack(
            anchor="w", padx=20, pady=(16, 10))

        # 模型选择（替代原版本选择）
        rm = ctk.CTkFrame(card_params, fg_color="transparent")
        rm.pack(fill="x", padx=20, pady=(0, 6))
        self._mk_label(rm, "模型选择", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 10))
        self.combo_batch_model = self._mk_optionmenu(
            rm, self._model_list(), width=240,
            command=self._on_batch_model_changed)
        self.combo_batch_model.pack(side="left")
        self.combo_batch_model.set(self._current_model())
        self.lbl_batch_model_hint = self._mk_label(rm, "", size=10,
                                                   color=COLORS["text_secondary"])
        self.lbl_batch_model_hint.pack(side="left", padx=(10, 0))

        # 尺寸
        r1 = ctk.CTkFrame(card_params, fg_color="transparent")
        r1.pack(fill="x", padx=20, pady=(0, 6))
        self._mk_label(r1, "尺寸比例", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 10))
        self.combo_batch_size = self._mk_optionmenu(
            r1, list(SIZE_MAP.keys()), width=170,
            command=self._on_batch_size_changed)
        self.combo_batch_size.pack(side="left")
        lk = self.config.get("last_size_key", DEFAULT_SIZE_KEY)
        self.combo_batch_size.set(lk if lk in SIZE_MAP else DEFAULT_SIZE_KEY)
        self.lbl_batch_pixel = self._mk_label(
            r1, self._size_caption(self.combo_batch_size.get(), self.combo_batch_model.get()),
            size=11, color=COLORS["text_secondary"])
        self.lbl_batch_pixel.pack(side="left", padx=(10, 0))

        # 出图画质（gpt-image-2 quality 档位，批量页可临时切换）
        rq = ctk.CTkFrame(card_params, fg_color="transparent")
        rq.pack(fill="x", padx=20, pady=(0, 6))
        self._mk_label(rq, "出图画质", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 10))
        self.combo_batch_quality = self._mk_optionmenu(
            rq, QUALITY_OPTIONS, width=120,
            command=self._on_batch_quality_changed)
        self.combo_batch_quality.pack(side="left")
        self.combo_batch_quality.set(self._quality())
        self._mk_label(rq, "（auto 由模型自动决定；high 更慢更贵）", size=10,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(8, 0))

        # 文件名
        r2 = ctk.CTkFrame(card_params, fg_color="transparent")
        r2.pack(fill="x", padx=20, pady=(0, 6))
        self._mk_label(r2, "文件名", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 10))
        self.entry_batch_prefix = self._mk_entry(r2)
        self.entry_batch_prefix.configure(width=160)
        self.entry_batch_prefix.pack(side="left")
        self.entry_batch_prefix.insert(0, self.config.get("filename_prefix", "ai_image"))
        self._mk_label(r2, ".png（自动追加序号与时间戳）", size=10,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(8, 0))

        # 保存路径
        r3 = ctk.CTkFrame(card_params, fg_color="transparent")
        r3.pack(fill="x", padx=20, pady=(0, 14))
        self._mk_label(r3, "保存到", size=12,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 10))
        self.lbl_batch_out_dir = self._mk_label(
            r3, self.config.get("save_dir", "") or "（未设置，请到「配置」页设置）",
            size=12, color=COLORS["text_primary"])
        self.lbl_batch_out_dir.pack(side="left")

        # -- 生成按钮 --
        bf = ctk.CTkFrame(page, fg_color="transparent")
        bf.pack(fill="x", padx=2, pady=(6, 0))
        self.ctx_batch.btn_generate = ctk.CTkButton(
            bf, text="🚀  批量生成", height=48, corner_radius=12,
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            font=("Microsoft YaHei UI", 15, "bold"),
            command=self._start_batch)
        self.ctx_batch.btn_generate.pack(fill="x")

        # -- 进度 + 一览 --
        self._build_progress_and_gallery(page, self.ctx_batch,
                                         "✨ 填写提示词列表后点击「批量生成」，生成结果在此一览")

        self._refresh_batch_model_hint()
        self._update_batch_count()
        return page

    def _build_progress_and_gallery(self, page, ctx, placeholder_text):
        """为某个生成页构建进度条 + 缩略图一览网格，存入 ctx"""
        ctx.progress_frame = ctk.CTkFrame(page, fg_color="transparent")
        ctx.progress_frame.pack(fill="x", padx=2, pady=(14, 0))
        ctx.progress_bar = ctk.CTkProgressBar(
            ctx.progress_frame, height=6, corner_radius=3,
            fg_color=COLORS["surface"], progress_color=COLORS["accent"],
            mode="determinate")
        ctx.lbl_progress = self._mk_label(ctx.progress_frame, "", size=12,
                                          color=COLORS["accent"])
        ctx.progress_frame.pack_forget()

        prev_card = self._mk_card(page)
        prev_card.pack(fill="both", expand=True, pady=(10, 0), padx=2)

        ph2 = ctk.CTkFrame(prev_card, fg_color="transparent")
        ph2.pack(fill="x", padx=20, pady=(16, 8))
        self._mk_label(ph2, "🖼️  生成一览", size=14, weight="bold").pack(side="left")
        self._mk_label(ph2, "双击图片可打开原图", size=10,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(10, 0))
        self._mk_btn(ph2, "📂  打开目录", "ghost",
                     command=self._open_save_dir).pack(side="right")

        ctx.gallery = ctk.CTkScrollableFrame(
            prev_card, fg_color=COLORS["surface"], corner_radius=8,
            height=260, orientation="vertical")
        ctx.gallery.pack(fill="both", expand=True, padx=20, pady=(0, 16))

        ctx.gallery_photos = []
        ctx.gallery_col = 0
        ctx.gallery_row = 0
        ctx.placeholder = self._mk_label(
            ctx.gallery, placeholder_text, size=13, color=COLORS["text_secondary"])
        ctx.placeholder.pack(pady=40)

    # ================================================================
    #  批量提示词行管理
    # ================================================================
    def _add_batch_row(self, text=""):
        if len(self.batch_rows) >= BATCH_MAX:
            messagebox.showinfo("提示", f"批量出图最多 {BATCH_MAX} 条提示词")
            return
        row = ctk.CTkFrame(self.batch_rows_frame, fg_color="transparent")
        row.pack(fill="x", pady=3)

        idx_lbl = self._mk_label(row, "", size=12, color=COLORS["text_secondary"],
                                 width=28)
        idx_lbl.pack(side="left", padx=(0, 6))

        entry = ctk.CTkEntry(
            row, height=36, fg_color=COLORS["surface"],
            border_color=COLORS["divider"], border_width=1, corner_radius=8,
            text_color=COLORS["text_primary"],
            placeholder_text="输入一条提示词，将单独生成一张图...",
            placeholder_text_color=COLORS["text_secondary"],
            font=("Microsoft YaHei UI", 13))
        entry.pack(side="left", fill="x", expand=True)
        if text:
            entry.insert(0, text)

        rec = {"frame": row, "entry": entry, "idx_lbl": idx_lbl}
        del_btn = ctk.CTkButton(
            row, text="✕", width=36, height=36, corner_radius=8,
            fg_color=COLORS["surface"], hover_color=COLORS["error"],
            text_color=COLORS["text_secondary"],
            font=("Microsoft YaHei UI", 13),
            command=lambda r=rec: self._remove_batch_row(r))
        del_btn.pack(side="left", padx=(6, 0))

        self.batch_rows.append(rec)
        self._renumber_batch_rows()
        self._update_batch_count()

    def _remove_batch_row(self, rec):
        if rec in self.batch_rows:
            self.batch_rows.remove(rec)
        rec["frame"].destroy()
        self._renumber_batch_rows()
        self._update_batch_count()

    def _clear_batch_rows(self):
        for rec in list(self.batch_rows):
            rec["frame"].destroy()
        self.batch_rows.clear()
        # 保留一行空白方便继续填
        self._add_batch_row()

    def _renumber_batch_rows(self):
        for i, rec in enumerate(self.batch_rows, 1):
            rec["idx_lbl"].configure(text=f"{i:>2}.")

    def _get_batch_prompts(self):
        return [rec["entry"].get().strip() for rec in self.batch_rows
                if rec["entry"].get().strip()]

    def _update_batch_count(self):
        if not hasattr(self, "lbl_batch_count"):
            return
        n = len(self._get_batch_prompts())
        if n == 0:
            self.lbl_batch_count.configure(text="请在上方填写提示词",
                                           text_color=COLORS["text_secondary"])
        else:
            self.lbl_batch_count.configure(
                text=f"共 {n} 条有效提示词 = {n} 张待生成（上限 {BATCH_MAX}）",
                text_color=COLORS["accent"])

    def _import_batch_prompts(self):
        path = filedialog.askopenfilename(
            title="导入提示词文件",
            filetypes=[("文本文件", "*.txt"), ("所有文件", "*.*")])
        if not path:
            return
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
        except UnicodeDecodeError:
            with open(path, "r", encoding="gbk", errors="ignore") as f:
                content = f.read()
        except Exception as e:
            messagebox.showerror("导入失败", str(e))
            return
        lines = [ln.strip() for ln in content.splitlines() if ln.strip()]
        if not lines:
            messagebox.showwarning("提示", "文件中没有有效提示词")
            return

        # 先把已有空行填掉，再追加新行
        added = 0
        for ln in lines:
            if len(self.batch_rows) >= BATCH_MAX:
                break
            # 找一个空行复用
            empty = next((r for r in self.batch_rows if not r["entry"].get().strip()), None)
            if empty:
                empty["entry"].insert(0, ln)
            else:
                self._add_batch_row(ln)
            added += 1
        self._update_batch_count()
        self.log_info(f"已导入 {added} 条提示词到批量列表"
                      + ("（已达上限，部分未导入）" if added < len(lines) else ""))

    # ================================================================
    #  参考图管理（单次 / 批量 共用，按 ctx 区分）
    # ================================================================
    def _build_ref_card(self, page, ctx, note="（本地图片将自动转码上传）",
                        style_lock=False):
        """构建紧凑型参考图卡片，组件引用存入 ctx。
        style_lock=True 时额外提供「锁定参考图风格」开关（批量出图用）。"""
        ref_card = self._mk_card(page)
        ref_card.pack(fill="x", pady=(0, 10), padx=2)

        # 标题行 + 操作按钮全部挤在一行，节省高度
        head = ctk.CTkFrame(ref_card, fg_color="transparent")
        head.pack(fill="x", padx=16, pady=(12, 6))
        self._mk_label(head, f"📎 参考图 {note}", size=12,
                       weight="bold").pack(side="left")
        ctx.lbl_ref_count = self._mk_label(head, "", size=10, color=COLORS["accent"])
        ctx.lbl_ref_count.pack(side="right")
        ctk.CTkButton(
            head, text="📁 选择", height=26, width=58, corner_radius=6,
            fg_color=COLORS["surface"], hover_color=COLORS["card_border"],
            border_color=COLORS["divider"], border_width=1,
            font=("Microsoft YaHei UI", 11),
            command=lambda c=ctx: self._add_ref_images(c)).pack(side="right", padx=(0, 8))
        ctk.CTkButton(
            head, text="清除", height=26, width=46, corner_radius=6,
            fg_color="transparent", hover_color=COLORS["surface"],
            text_color=COLORS["text_secondary"],
            font=("Microsoft YaHei UI", 11),
            command=lambda c=ctx: self._clear_ref_images(c)).pack(side="right", padx=(0, 6))

        # 缩略图行（小尺寸）—— 固定较矮高度，空状态不再撑出大片留白
        ctx.ref_thumb_frame = ctk.CTkFrame(ref_card, fg_color="transparent",
                                            height=58)
        ctx.ref_thumb_frame.pack_propagate(False)
        ctx.ref_thumb_frame.pack(fill="x", padx=16, pady=(0, 4))

        # URL 行（紧凑）—— gpt-image-2 图生图走文件上传，URL 仅作占位提示
        url_row = ctk.CTkFrame(ref_card, fg_color="transparent")
        url_row.pack(fill="x", padx=16, pady=(0, 12))
        self._mk_label(url_row, "或URL", size=10,
                       color=COLORS["text_secondary"]).pack(side="left", padx=(0, 6))
        ctx.entry_ref_urls = ctk.CTkEntry(
            url_row, height=28, fg_color=COLORS["surface"],
            border_color=COLORS["divider"], border_width=1, corner_radius=6,
            placeholder_text="当前模型图生图请用「选择」上传图片，URL 暂不支持",
            placeholder_text_color=COLORS["text_secondary"],
            font=("Consolas", 10))
        ctx.entry_ref_urls.pack(side="left", fill="x", expand=True)

        # 风格锁定开关（仅批量页）：勾选后给每条提示词自动注入"沿用参考图风格"指令
        if style_lock:
            sl_row = ctk.CTkFrame(ref_card, fg_color="transparent")
            sl_row.pack(fill="x", padx=16, pady=(0, 12))
            ctx.style_lock_var = ctk.BooleanVar(value=True)
            ctk.CTkCheckBox(
                sl_row, text="  锁定参考图风格（让每张图都沿用参考图的配色/版式/质感，保持系列统一）",
                variable=ctx.style_lock_var,
                checkbox_width=18, checkbox_height=18, corner_radius=5,
                fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
                border_color=COLORS["divider"], border_width=2,
                text_color=COLORS["text_primary"],
                font=("Microsoft YaHei UI", 11)).pack(side="left")

    def _add_ref_images(self, ctx):
        files = filedialog.askopenfilenames(
            title="选择参考图片",
            filetypes=[("图片文件", "*.png *.jpg *.jpeg *.webp *.bmp"), ("所有文件", "*.*")]
        )
        if not files:
            return
        for f in files:
            if f not in ctx.ref_paths and len(ctx.ref_paths) < 16:
                ctx.ref_paths.append(f)
                self._add_ref_thumbnail(ctx, f)
        self._update_ref_count(ctx)
        self.log_info(f"已添加参考图，共 {len(ctx.ref_paths)} 张")

    def _add_ref_thumbnail(self, ctx, path):
        try:
            img = Image.open(path)
            img.thumbnail((44, 44), Image.LANCZOS)
            photo = ImageTk.PhotoImage(img)
            ctx.ref_photos.append(photo)

            container = ctk.CTkFrame(ctx.ref_thumb_frame, fg_color=COLORS["surface"],
                                     corner_radius=6, width=52, height=52)
            container.pack_propagate(False)
            container.pack(side="left", padx=3, pady=3)

            lbl = ctk.CTkLabel(container, image=photo, text="", cursor="hand2")
            lbl.place(relx=0.5, rely=0.5, anchor="center")

            def _remove(p=path, c=container, ph=photo):
                if p in ctx.ref_paths:
                    ctx.ref_paths.remove(p)
                if ph in ctx.ref_photos:
                    ctx.ref_photos.remove(ph)
                c.destroy()
                self._update_ref_count(ctx)

            lbl.bind("<Button-1>", lambda e: _remove())
        except Exception as e:
            self.log_error(f"加载缩略图失败: {e}")

    def _clear_ref_images(self, ctx):
        ctx.ref_paths.clear()
        ctx.ref_photos.clear()
        for w in ctx.ref_thumb_frame.winfo_children():
            w.destroy()
        self._update_ref_count(ctx)

    def _update_ref_count(self, ctx):
        n = len(ctx.ref_paths)
        ctx.lbl_ref_count.configure(text=f"{n} 张（点击缩略图可删除）" if n > 0 else "")

    def _get_ref_files(self, ctx):
        """把 ctx 的本地参考图压缩为字节流，供 multipart 上传 /v1/images/edits。
        返回 [(filename, bytes, mime), ...]，最多 16 张。单张压到 ~1.5MB 以内最稳。"""
        result = []
        for p in ctx.ref_paths[:16]:
            try:
                img = Image.open(p)
                # 长边超 2048 先缩小，控制体积
                max_side = max(img.width, img.height)
                if max_side > 2048:
                    ratio = 2048 / max_side
                    img = img.resize((int(img.width * ratio), int(img.height * ratio)),
                                     Image.LANCZOS)
                buf = BytesIO()
                quality = 90
                img.convert("RGB").save(buf, format="JPEG", quality=quality)
                # 逐步压到 1.5MB 以内
                while buf.tell() > 1536 * 1024 and quality > 30:
                    quality -= 10
                    buf.seek(0)
                    buf.truncate()
                    img.convert("RGB").save(buf, format="JPEG", quality=quality)
                fbytes = buf.getvalue()
                fname = os.path.splitext(os.path.basename(p))[0] + ".jpg"
                result.append((fname, fbytes, "image/jpeg"))
                self.log_info(f"参考图 {os.path.basename(p)}: "
                              f"{img.width}x{img.height} → {len(fbytes)/1024:.0f}KB")
            except Exception as e:
                self.log_error(f"参考图失败: {os.path.basename(p)} — {e}")
        return result

    def _warn_pasted_urls(self, ctx):
        """gpt-image-2 图生图走 multipart 文件上传，不支持外部 URL。
        若用户在 URL 框粘了内容，给出提示但不使用。"""
        if not getattr(ctx, "entry_ref_urls", None):
            return
        text = ctx.entry_ref_urls.get().strip()
        if text:
            self.log_info("提示: 当前模型图生图通过上传图片文件实现，"
                          "URL 粘贴框对其无效，已忽略。请用「选择」上传本地图片。")

    # ================================================================
    #  导航 & 交互
    # ================================================================
    def _switch_tab(self, tab):
        if self.current_page == tab:
            return
        self.page_config.pack_forget()
        self.page_analyze.pack_forget()
        self.page_single.pack_forget()
        self.page_batch.pack_forget()
        if tab == "config":
            self.page_config.pack(fill="both", expand=True, pady=(4, 0))
        elif tab == "analyze":
            self.page_analyze.pack(fill="both", expand=True, pady=(4, 0))
        elif tab == "single":
            self.page_single.pack(fill="both", expand=True, pady=(4, 0))
        else:
            self.page_batch.pack(fill="both", expand=True, pady=(4, 0))
        self.current_page = tab
        self._highlight_nav(tab)

    def _highlight_nav(self, active):
        for btn, tab in [(self.nav_config, "config"),
                         (self.nav_analyze, "analyze"),
                         (self.nav_single, "single"),
                         (self.nav_batch, "batch")]:
            if tab == active:
                btn.configure(fg_color=COLORS["accent"], text_color="white",
                              hover_color=COLORS["accent_hover"])
            else:
                btn.configure(fg_color="transparent", text_color=COLORS["text_secondary"],
                              hover_color=COLORS["surface"])

    def _toggle_key_vis(self):
        if self.entry_key.cget("show") == "•":
            self.entry_key.configure(show="")
            self.btn_toggle_key.configure(text="🙈")
        else:
            self.entry_key.configure(show="•")
            self.btn_toggle_key.configure(text="👁")

    @staticmethod
    def _open_url(url):
        """用系统默认浏览器打开链接。"""
        import webbrowser
        webbrowser.open(url)

    def _save_config(self):
        self.config["api_key"] = self.entry_key.get().strip()
        self.config["api_base"] = self.entry_base.get().strip()
        self.config["save_dir"] = self.entry_save_dir.get().strip()
        save_config(self.config)
        self._check_api_status()
        self._refresh_save_dir_labels()
        self._set_status("✅ 配置已保存", COLORS["success"])
        self.log_success("配置已保存")
        self.win.after(2000, lambda: self._set_status("就绪", COLORS["text_secondary"]))

    def _browse_save_dir(self):
        d = filedialog.askdirectory(title="选择保存目录")
        if d:
            self.entry_save_dir.delete(0, "end")
            self.entry_save_dir.insert(0, d)
            self.config["save_dir"] = d
            save_config(self.config)
            self._refresh_save_dir_labels()

    def _refresh_save_dir_labels(self):
        d = self.config.get("save_dir", "") or "（未设置，请到「配置」页设置）"
        if hasattr(self, "lbl_single_out_dir"):
            self.lbl_single_out_dir.configure(text=d)
        if hasattr(self, "lbl_batch_out_dir"):
            self.lbl_batch_out_dir.configure(text=d)

    def _size_caption(self, size_key, model_name=None):
        info = SIZE_MAP.get(size_key, {"pixel": "1024x1024", "ratio": "1:1"})
        px = info["pixel"]
        return "(由模型自动决定)" if px == "auto" else f"({px})"

    # ================================================================
    #  模型管理（查询 / 手动 / 多保存）
    # ================================================================
    def _model_list(self):
        """返回用户已保存的模型列表（至少含一个占位，避免下拉为空）。"""
        models = self.config.get("models") or []
        return models if models else ["gpt-image-2"]

    def _current_model(self):
        models = self._model_list()
        cur = self.config.get("model", "")
        return cur if cur in models else models[0]

    def _refresh_saved_models(self):
        """重绘配置页「已保存模型」列表。"""
        if not hasattr(self, "frame_saved_models"):
            return
        for w in self.frame_saved_models.winfo_children():
            w.destroy()
        models = self.config.get("models") or []
        if not models:
            self._mk_label(self.frame_saved_models, "（暂无，请查询或手动添加）",
                           size=11, color=COLORS["text_secondary"]).pack(
                anchor="w", padx=12, pady=8)
            return
        for m in models:
            row = ctk.CTkFrame(self.frame_saved_models, fg_color="transparent")
            row.pack(fill="x", padx=10, pady=3)
            self._mk_label(row, f"• {m}", size=12).pack(side="left")
            ctk.CTkButton(
                row, text="删除", height=24, width=48, corner_radius=6,
                fg_color="transparent", hover_color=COLORS["error"],
                text_color=COLORS["text_secondary"],
                font=("Microsoft YaHei UI", 11),
                command=lambda name=m: self._remove_model(name)).pack(side="right")

    def _sync_model_dropdowns(self):
        """模型列表变动后，同步更新单次/批量页的模型下拉。"""
        values = self._model_list()
        cur = self._current_model()
        for combo in ("combo_single_model", "combo_batch_model"):
            if hasattr(self, combo):
                getattr(self, combo).configure(values=values)
                getattr(self, combo).set(cur)
        # 文本对话页：同样保持与已保存模型一致；尽量保留当前选择
        if hasattr(self, "combo_chat_model"):
            self.combo_chat_model.configure(values=values)
            cur_chat = self.combo_chat_model.get()
            if cur_chat not in values:
                cur_chat = values[0] if values else "gpt-4o"
                self.combo_chat_model.set(cur_chat)
                self._on_chat_model_changed(cur_chat)

    def _add_model(self, name):
        name = (name or "").strip()
        if not name:
            return
        models = self.config.get("models") or []
        if name in models:
            self._set_status(f"模型 {name} 已存在", COLORS["warning"])
            return
        if len(models) >= MODEL_SAVE_LIMIT:
            messagebox.showinfo(
                "提示",
                f"最多只能保存 {MODEL_SAVE_LIMIT} 个模型。\n"
                f"请先在「配置」页删除部分已保存模型后再添加。")
            self._set_status(
                f"已达上限 {MODEL_SAVE_LIMIT} 个，未添加 {name}", COLORS["warning"])
            return
        models.append(name)
        self.config["models"] = models
        self.config["model"] = name
        save_config(self.config)
        self._refresh_saved_models()
        self._sync_model_dropdowns()
        self.log_success(f"已添加模型: {name}")
        self._set_status(f"✅ 已添加模型 {name}", COLORS["success"])

    def _add_manual_model(self):
        name = self.entry_manual_model.get().strip()
        if not name:
            messagebox.showinfo("提示", "请输入模型 ID")
            return
        self._add_model(name)
        self.entry_manual_model.delete(0, "end")

    def _remove_model(self, name):
        models = self.config.get("models") or []
        if name in models:
            models.remove(name)
        self.config["models"] = models
        if self.config.get("model") == name:
            self.config["model"] = models[0] if models else "gpt-image-2"
        save_config(self.config)
        self._refresh_saved_models()
        self._sync_model_dropdowns()
        self.log_info(f"已删除模型: {name}")

    def _query_models(self):
        """查询云雾可用模型（GET /v1/models），子线程执行，失败可重试。"""
        api_key = self.entry_key.get().strip() or self.config.get("api_key", "")
        api_base = self.entry_base.get().strip() or self.config.get("api_base", "")
        if not api_key:
            messagebox.showwarning("提示", "请先填入 API Key 再查询模型")
            return
        self.btn_query_models.configure(state="disabled", text="查询中...")
        self.lbl_query_hint.configure(text="正在请求 /v1/models ...",
                                      text_color=COLORS["text_secondary"])

        def worker():
            try:
                ids, info = openai_list_models(api_base, api_key)
                self._ui(lambda: self._on_models_queried(ids))
            except Exception as e:
                msg = str(e)
                self._ui(lambda m=msg: self._on_models_query_failed(m))

        threading.Thread(target=worker, daemon=True).start()

    def _on_models_queried(self, ids):
        self.btn_query_models.configure(state="normal", text="🔍 查询可用模型")
        if not ids:
            self.lbl_query_hint.configure(text="未返回模型，可手动输入",
                                          text_color=COLORS["warning"])
            return
        self.lbl_query_hint.configure(text=f"查到 {len(ids)} 个模型",
                                      text_color=COLORS["success"])
        self.log_success(f"查询到 {len(ids)} 个可用模型")
        self._show_model_picker(ids)

    def _on_models_query_failed(self, msg):
        self.btn_query_models.configure(state="normal", text="🔄 重试查询")
        self.lbl_query_hint.configure(text=f"查询失败：{msg[:50]}（可重试或手动输入）",
                                      text_color=COLORS["error"])
        self.log_error(f"查询模型失败: {msg[:120]}")

    def _show_model_picker(self, ids):
        """弹窗列出查询到的模型，支持模糊搜索过滤 + 多选勾选，确认后批量加入。
        性能优化：搜索框防抖（250ms）+ 渲染上限（MODEL_RENDER_LIMIT 条），
        避免一次性渲染数百个复选框导致界面卡死。
        约束：保存的模型总数最多 MODEL_SAVE_LIMIT 个。"""
        win = ctk.CTkToplevel(self.win)
        win.title("选择要添加的模型")
        win.geometry("460x600")
        win.configure(fg_color=COLORS["bg"])
        win.transient(self.win)
        win.grab_set()

        existing = set(self.config.get("models") or [])
        existing_n = len(existing)            # 已保存数（计入 10 个上限）
        # 勾选状态用字典持久保存（跨过滤刷新不丢失），key=模型ID
        checked = {mid: False for mid in ids}
        # 防抖句柄
        debounce = {"id": None}

        # —— 预归一化：把每个模型 ID 转成「全小写、去除所有分隔符」的形式，
        #    例如 "Claude-3-Opus" -> "claude3opus"，这样无论用户搜
        #    "claude opus" / "claude-opus" / "claudeopus" 都能命中。
        import re as _re
        _norm_cache = {mid: _re.sub(r"[\s\-_./:]+", "", mid.lower()) for mid in ids}

        def _match(mid, tokens):
            """多关键词 AND 匹配：每个关键词都要在归一化后的模型名里出现。
            兼顾原始名与归一化名，确保连字符/空格/大小写都不影响命中。"""
            raw = mid.lower()
            norm = _norm_cache[mid]
            for t in tokens:
                if t not in raw and t not in norm:
                    return False
            return True

        self._mk_label(win, f"云雾 API 可用模型（{len(ids)} 个）", size=14,
                       weight="bold").pack(anchor="w", padx=20, pady=(16, 4))
        self._mk_label(
            win, f"输入关键词模糊筛选，勾选后点底部「添加所选」"
                 f"（最多保存 {MODEL_SAVE_LIMIT} 个）",
            size=11, color=COLORS["text_secondary"]).pack(
            anchor="w", padx=20, pady=(0, 8))

        # 搜索框 + 清空
        search_row = ctk.CTkFrame(win, fg_color="transparent")
        search_row.pack(fill="x", padx=20, pady=(0, 6))
        entry_search = self._mk_entry(search_row)
        entry_search.configure(
            placeholder_text="🔍 关键词过滤，支持空格分词，如：claude opus / gpt image")
        entry_search.pack(side="left", fill="x", expand=True)

        # 匹配计数（搜索框正下方，实时反馈）
        lbl_match = self._mk_label(win, "", size=11, color=COLORS["text_secondary"])
        lbl_match.pack(anchor="w", padx=20, pady=(0, 4))

        scroll = ctk.CTkScrollableFrame(win, fg_color=COLORS["surface"], corner_radius=8)
        scroll.pack(fill="both", expand=True, padx=20, pady=(6, 6))

        lbl_count = self._mk_label(win, "", size=11, color=COLORS["accent"])
        lbl_count.pack(anchor="w", padx=20, pady=(0, 4))

        def _saved_total():
            """当前会勾选保存后的总数 = 已保存 + 本次新勾选(不在已保存里的)。"""
            new_checked = sum(1 for m, ok in checked.items()
                              if ok and m not in existing)
            return existing_n + new_checked

        def render(keyword=""):
            for w in scroll.winfo_children():
                w.destroy()
            # 关键词拆分：按空格切成多个 token，每个都归一化（去分隔符、转小写），
            # 实现「多关键词 AND 匹配」。如搜 "claude opus" 会拆成 ["claude","opus"]。
            raw_kw = keyword.strip().lower()
            tokens = [_re.sub(r"[\s\-_./:]+", "", t)
                      for t in raw_kw.split() if t.strip()]
            tokens = [t for t in tokens if t]
            if tokens:
                shown = [m for m in ids if _match(m, tokens)]
            else:
                shown = list(ids)
            if not shown:
                self._mk_label(
                    scroll,
                    f"未找到包含「{keyword.strip()}」的模型\n"
                    f"· 可尝试更短的关键词，或用空格分词（如：claude opus）\n"
                    f"· 当前共 {len(ids)} 个可选模型",
                    size=12, color=COLORS["text_secondary"],
                    justify="left").pack(anchor="w", padx=12, pady=10)
                lbl_match.configure(text=f"共匹配 0 个 / {len(ids)} 个模型")
                return
            # 渲染上限：超出只显示前 N 条，提示用户细化关键词
            truncated = len(shown) > MODEL_RENDER_LIMIT
            rendered = min(len(shown), MODEL_RENDER_LIMIT)
            for mid in shown[:MODEL_RENDER_LIMIT]:
                v = ctk.BooleanVar(value=checked[mid])

                def _on_toggle(name=mid, var=v):
                    want = var.get()
                    # 勾选时校验 10 个上限（取消勾选无需校验）
                    if want and name not in existing and _saved_total() >= MODEL_SAVE_LIMIT:
                        var.set(False)          # 撤销本次勾选
                        checked[name] = False
                        messagebox.showinfo(
                            "提示",
                            f"最多只能保存 {MODEL_SAVE_LIMIT} 个模型。\n"
                            f"如需添加新模型，请先在「配置」页删除部分已保存模型。")
                        return
                    checked[name] = want
                    _update_count()

                label = f"  {mid}" + ("  (已保存)" if mid in existing else "")
                ctk.CTkCheckBox(
                    scroll, text=label, variable=v, command=_on_toggle,
                    checkbox_width=18, checkbox_height=18, corner_radius=5,
                    fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
                    border_color=COLORS["divider"], border_width=2,
                    text_color=COLORS["text_primary"],
                    font=("Microsoft YaHei UI", 12)).pack(anchor="w", padx=12, pady=4)
            if truncated:
                self._mk_label(
                    scroll,
                    f"… 共匹配 {len(shown)} 个，仅显示前 {MODEL_RENDER_LIMIT} 个，"
                    f"请输入更精确的关键词缩小范围",
                    size=11, color=COLORS["warning"]).pack(
                    anchor="w", padx=12, pady=(8, 6))
            # 匹配计数标签：实时反馈匹配/渲染情况
            if tokens:
                lbl_match.configure(
                    text=f"共匹配 {len(shown)} 个 / {len(ids)} 个，已显示 {rendered} 个")
            else:
                lbl_match.configure(
                    text=f"共 {len(ids)} 个模型，已显示 {rendered} 个")

        def _update_count():
            n = sum(1 for ok in checked.values() if ok)
            lbl_count.configure(
                text=f"已勾选 {n} 个  |  保存后共 {_saved_total()}/{MODEL_SAVE_LIMIT} 个")

        def _clear_all():
            for m in checked:
                checked[m] = False
            render(entry_search.get())
            _update_count()

        self._mk_btn(search_row, "清空勾选", "ghost",
                     command=_clear_all).pack(side="left", padx=(8, 0))

        def _on_key(_e=None):
            # 防抖：取消上一个待执行的渲染，250ms 后再渲染，避免每次按键全量重绘
            if debounce["id"] is not None:
                try:
                    win.after_cancel(debounce["id"])
                except Exception:
                    pass
            debounce["id"] = win.after(
                250, lambda: render(entry_search.get()))

        entry_search.bind("<KeyRelease>", _on_key)

        def do_add():
            chosen = [m for m, ok in checked.items() if ok]
            added = 0
            for m in chosen:
                if m not in (self.config.get("models") or []):
                    self._add_model(m)
                    added += 1
            win.destroy()
            if added:
                self._set_status(f"✅ 已添加 {added} 个模型", COLORS["success"])

        bf = ctk.CTkFrame(win, fg_color="transparent")
        bf.pack(fill="x", padx=20, pady=(0, 16))
        self._mk_btn(bf, "添加所选", "primary", command=do_add).pack(side="left")
        self._mk_btn(bf, "取消", "secondary", command=win.destroy).pack(side="left", padx=(8, 0))

        render()
        _update_count()

    def _on_single_quality_changed(self, choice):
        """单次页画质切换：持久化并同步到批量页。"""
        self.config["quality"] = choice
        save_config(self.config)
        if hasattr(self, "combo_batch_quality"):
            self.combo_batch_quality.set(choice)

    def _on_batch_quality_changed(self, choice):
        """批量页画质切换：持久化并同步到单次页。"""
        self.config["quality"] = choice
        save_config(self.config)
        if hasattr(self, "combo_single_quality"):
            self.combo_single_quality.set(choice)

    # ---- 单次页模型/尺寸 ----
    def _on_single_model_changed(self, choice):
        self.config["model"] = choice
        save_config(self.config)
        if hasattr(self, "combo_batch_model"):
            self.combo_batch_model.set(choice)
        self.lbl_single_pixel.configure(
            text=self._size_caption(self.combo_single_size.get(), choice))

    def _on_single_size_changed(self, choice):
        self.lbl_single_pixel.configure(text=self._size_caption(choice))
        self.config["last_size_key"] = choice
        save_config(self.config)
        if hasattr(self, "combo_batch_size"):
            self.combo_batch_size.set(choice)
            self.lbl_batch_pixel.configure(text=self._size_caption(choice))

    def _refresh_single_model_hint(self):
        if not hasattr(self, "lbl_single_model_hint"):
            return
        self.lbl_single_model_hint.configure(
            text=f"画质 {self.config.get('quality', 'auto')}")

    # ---- 批量页模型/尺寸 ----
    def _on_batch_model_changed(self, choice):
        self.config["model"] = choice
        save_config(self.config)
        if hasattr(self, "combo_single_model"):
            self.combo_single_model.set(choice)
        self.lbl_batch_pixel.configure(
            text=self._size_caption(self.combo_batch_size.get(), choice))

    def _on_batch_size_changed(self, choice):
        self.lbl_batch_pixel.configure(text=self._size_caption(choice))
        self.config["last_size_key"] = choice
        save_config(self.config)
        if hasattr(self, "combo_single_size"):
            self.combo_single_size.set(choice)
            self.lbl_single_pixel.configure(text=self._size_caption(choice))

    def _refresh_batch_model_hint(self):
        if not hasattr(self, "lbl_batch_model_hint"):
            return
        self.lbl_batch_model_hint.configure(
            text=f"画质 {self.config.get('quality', 'auto')}")

    def _check_api_status(self):
        has = bool(self.entry_key.get().strip() or self.config.get("api_key", ""))
        self.lbl_api_status.configure(
            text="🟢 API 已配置" if has else "🔴 未配置 API",
            text_color=COLORS["success"] if has else COLORS["error"])

    # ================================================================
    #  一览网格
    # ================================================================
    def _reset_gallery(self, ctx):
        for w in ctx.gallery.winfo_children():
            w.destroy()
        ctx.gallery_photos = []
        ctx.gallery_col = 0
        ctx.gallery_row = 0

    def _add_gallery_item(self, ctx, filepath, ok=True, caption=""):
        cols = 3
        cell = ctk.CTkFrame(ctx.gallery, fg_color=COLORS["card"],
                            corner_radius=8, width=150, height=170)
        cell.grid(row=ctx.gallery_row, column=ctx.gallery_col,
                  padx=6, pady=6, sticky="n")
        cell.grid_propagate(False)

        if ok:
            try:
                img = Image.open(filepath)
                img.thumbnail((130, 120), Image.LANCZOS)
                photo = ImageTk.PhotoImage(img)
                ctx.gallery_photos.append(photo)
                lbl = ctk.CTkLabel(cell, image=photo, text="", cursor="hand2")
                lbl.pack(pady=(8, 4))
                lbl.bind("<Double-Button-1>",
                         lambda e, p=filepath: self._open_file(p))
                ctk.CTkLabel(cell, text=os.path.basename(filepath),
                             font=("Microsoft YaHei UI", 9),
                             text_color=COLORS["text_secondary"],
                             wraplength=134).pack(pady=(0, 6))
            except Exception:
                ctk.CTkLabel(cell, text="✅ 已保存\n(预览失败)",
                             font=("Microsoft YaHei UI", 11),
                             text_color=COLORS["success"]).pack(expand=True)
        else:
            ctk.CTkLabel(cell, text="❌", font=("Microsoft YaHei UI", 28),
                         text_color=COLORS["error"]).pack(pady=(20, 4))
            ctk.CTkLabel(cell, text=caption or "失败",
                         font=("Microsoft YaHei UI", 9),
                         text_color=COLORS["error"],
                         wraplength=134).pack(pady=(0, 6))

        ctx.gallery_col += 1
        if ctx.gallery_col >= cols:
            ctx.gallery_col = 0
            ctx.gallery_row += 1

    def _set_status(self, msg, color=None):
        self.status_bar.configure(text=msg, text_color=color or COLORS["text_secondary"])

    def _open_save_dir(self):
        d = self.config.get("save_dir", "")
        if d and os.path.isdir(d):
            self._open_file(d)
        else:
            messagebox.showwarning("提示", "保存目录无效，请到「配置」页设置")

    # ================================================================
    #  通用前置校验
    # ================================================================
    def _validate_common(self):
        """返回 (api_key, api_base, save_dir) 或 None。"""
        api_key = self.entry_key.get().strip() or self.config.get("api_key", "")
        if not api_key:
            messagebox.showwarning("提示", "请先在「配置」页填入 API Key")
            self._switch_tab("config")
            return None
        api_base = self.entry_base.get().strip() or self.config["api_base"]
        save_dir = self.config.get("save_dir", "").strip()
        if not save_dir or not os.path.isdir(save_dir):
            messagebox.showwarning(
                "提示", f"默认保存目录无效或不存在，请到「配置」页设置:\n{save_dir}")
            self._switch_tab("config")
            return None
        return api_key, api_base, save_dir

    def _quality(self):
        """当前画质档位（gpt-image-2 的 quality 字段），默认 auto。
        只接受 QUALITY_OPTIONS 内的值，旧配置（如速创时代的 4K）一律回退 auto，
        避免把非法画质传给云雾 API 导致调用失败。"""
        q = self.config.get("quality", "auto")
        return q if q in QUALITY_OPTIONS else "auto"

    # ================================================================
    #  单次出图流程
    # ================================================================
    def _start_single(self):
        if self.generating:
            return
        prompt = self.text_single_prompt.get("1.0", "end-1c").strip()
        if not prompt:
            messagebox.showwarning("提示", "请输入提示词")
            return
        model_name = self.combo_single_model.get()
        common = self._validate_common()
        if not common:
            return
        api_key, api_base, save_dir = common

        size_key = self.combo_single_size.get()
        prefix = self.entry_single_prefix.get().strip() or "ai_image"
        quality = self._quality()

        self._warn_pasted_urls(self.ctx_single)
        ref_files = self._get_ref_files(self.ctx_single)

        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        jobs = [{"seq": 1, "prompt": prompt,
                 "filepath": os.path.join(save_dir, f"{prefix}_{ts}.png")}]

        self._persist_gen_config(api_key, api_base, save_dir, prefix, size_key)

        if not self.log_expanded:
            self._toggle_log()
        mode = "图生图" if ref_files else "文生图"
        self.log_info(f"单次开始 | 模型={model_name} | {mode} | 尺寸={size_key} "
                      f"| 画质={quality} | 参考图={len(ref_files)}")

        self._begin_generation(self.ctx_single, "⏳  生成中...",
                               api_base, api_key, size_key, model_name,
                               quality, ref_files, jobs)

    # ================================================================
    #  批量出图流程
    # ================================================================
    def _start_batch(self):
        if self.generating:
            return
        prompts = self._get_batch_prompts()
        if not prompts:
            messagebox.showwarning("提示", "请在提示词列表中至少填写一条")
            return
        if len(prompts) > BATCH_MAX:
            messagebox.showwarning("提示", f"批量出图最多 {BATCH_MAX} 张")
            return
        model_name = self.combo_batch_model.get()
        common = self._validate_common()
        if not common:
            return
        api_key, api_base, save_dir = common

        size_key = self.combo_batch_size.get()
        prefix = self.entry_batch_prefix.get().strip() or "ai_image"
        quality = self._quality()

        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        jobs = []
        for i, prompt in enumerate(prompts, 1):
            fname = f"{prefix}_{ts}_{i:02d}.png"
            jobs.append({"seq": i, "prompt": prompt,
                         "filepath": os.path.join(save_dir, fname)})
        total = len(jobs)

        if total > 20:
            if not messagebox.askyesno(
                    "确认批量生成",
                    f"本次将生成 {total} 张图片（每条提示词 1 张）。\n按张计费，确认继续？"):
                return

        # 批量页参考图：作为模板应用到所有提示词
        self._warn_pasted_urls(self.ctx_batch)
        ref_files = self._get_ref_files(self.ctx_batch)

        # 风格锁定：勾选且有参考图时，给每条提示词注入"沿用参考图风格"前缀
        style_lock = bool(self.ctx_batch.style_lock_var
                          and self.ctx_batch.style_lock_var.get())
        if style_lock and ref_files:
            for job in jobs:
                job["prompt"] = STYLE_LOCK_DIRECTIVE + job["prompt"]
            self.log_info("已启用「锁定参考图风格」：每条提示词均注入统一风格指令")
        elif style_lock and not ref_files:
            self.log_info("提示: 已勾选锁定风格，但未提供参考图，风格指令不生效")

        self._persist_gen_config(api_key, api_base, save_dir, prefix, size_key)

        if not self.log_expanded:
            self._toggle_log()
        mode = "图生图" if ref_files else "文生图"
        self.log_info(f"批量开始 | 模型={model_name} | {mode} | 尺寸={size_key} "
                      f"| 画质={quality} | {total} 条提示词 = {total} 张 "
                      f"| 参考图={len(ref_files)} "
                      f"| 并发={min(MAX_CONCURRENCY, total)}")

        self._begin_generation(self.ctx_batch, "⏳  批量生成中...",
                               api_base, api_key, size_key, model_name,
                               quality, ref_files, jobs)

    def _persist_gen_config(self, api_key, api_base, save_dir, prefix, size_key):
        for key, val in [("api_key", api_key), ("api_base", api_base),
                         ("save_dir", save_dir), ("filename_prefix", prefix),
                         ("last_size_key", size_key)]:
            self.config[key] = val
        save_config(self.config)

    # ================================================================
    #  通用生成执行（两页共用）
    # ================================================================
    def _begin_generation(self, ctx, busy_text, api_base, api_key, size_key,
                          model_name, quality, ref_files, jobs):
        total = len(jobs)
        self.generating = True
        ctx.btn_generate.configure(text=busy_text, state="disabled",
                                   fg_color=COLORS["surface"],
                                   text_color=COLORS["text_secondary"])
        # 锁住另一个页的生成按钮，避免并行冲突
        other = self.ctx_batch if ctx is self.ctx_single else self.ctx_single
        other.btn_generate.configure(state="disabled")

        ctx.progress_frame.pack(fill="x", padx=2, pady=(14, 0))
        ctx.progress_bar.configure(mode="determinate")
        ctx.progress_bar.set(0)
        ctx.progress_bar.pack(fill="x")
        ctx.lbl_progress.pack(anchor="w", pady=(4, 0))
        ctx.lbl_progress.configure(text=f"0 / {total} 已完成")
        self._reset_gallery(ctx)

        t = threading.Thread(
            target=self._batch_worker,
            args=(ctx, api_base, api_key, size_key, model_name, quality, ref_files, jobs),
            daemon=True)
        t.start()

    def _batch_worker(self, ctx, api_base, api_key, size_key, model_name,
                      quality, ref_files, jobs):
        total = len(jobs)
        done = {"n": 0, "ok": 0, "fail": 0}
        lock = threading.Lock()

        def run_one(job):
            try:
                fp = self._generate_one(api_base, api_key, job["prompt"],
                                        size_key, model_name, quality,
                                        job["filepath"], ref_files, job["seq"])
                with lock:
                    done["ok"] += 1
                self._ui(lambda p=fp: self._add_gallery_item(ctx, p, ok=True))
            except Exception as e:
                msg = str(e)
                with lock:
                    done["fail"] += 1
                self.log_error(f"#{job['seq']} 失败: {msg[:80]}")
                cap = msg[:40]
                self._ui(lambda c=cap: self._add_gallery_item(ctx, "", ok=False, caption=c))
            finally:
                with lock:
                    done["n"] += 1
                    n = done["n"]
                frac = n / total
                self._ui(lambda f=frac, c=n: (
                    ctx.progress_bar.set(f),
                    ctx.lbl_progress.configure(text=f"{c} / {total} 已完成")))

        workers = min(MAX_CONCURRENCY, total)
        with ThreadPoolExecutor(max_workers=workers) as ex:
            futures = [ex.submit(run_one, j) for j in jobs]
            for _ in as_completed(futures):
                pass

        self._ui(lambda: self._batch_done(ctx, done["ok"], done["fail"], total))

    @staticmethod
    def _desensitize_amount(text):
        """对提示词中的金额数字进行脱敏，替换为 xxx，保留金额单位。

        匹配规则（按优先级）：
        1. 货币符号 + 数字        → $5000 / ¥1,200 → xxx
        2. 数字 + 带量纲单位       → 5000万元 / 3.5亿元 / 1200元 → xxx万元 / xxx亿元 / xxx元
        3. 数字 + 裸单位（万/亿）  → 5000万的业绩 → xxx万的业绩
        """
        # 规则1：货币符号 + 数字（含千分位逗号），符号和数字一起替换
        text = re.sub(r'[$¥€£￥]\s*[\d,]+\.?\d*', 'xxx', text)

        # 规则2：数字 + 带量纲单位，只替换数字，保留单位
        text = re.sub(
            r'([\d,]+\.?\d*)\s*(百万元|万元|亿元|千万元|十亿元|万块|亿块|百万块|千元|百元|十元|元|块|角|分|毛)',
            r'xxx\2', text)

        # 规则3：数字 + 裸单位（千万/万/亿），只替换数字，保留单位
        text = re.sub(
            r'([\d,]+\.?\d*)\s*(千万|万千?|亿万?)',
            r'xxx\2', text)

        return text

    def _generate_one(self, api_base, api_key, prompt, size_key, model_name,
                      quality, filepath, ref_files, seq):
        """OpenAI 兼容同步出图：有参考图走图生图(edits)，否则文生图(generations)。"""
        # ── 金额脱敏 ──
        safe_prompt = self._desensitize_amount(prompt)
        if safe_prompt != prompt:
            self.log_info(f"#{seq} 检测到金额数字，已脱敏处理")

        size_info = SIZE_MAP.get(size_key, {"pixel": "1024x1024", "ratio": "1:1"})
        size_pixel = size_info["pixel"]

        if ref_files:
            img_data, log_data = openai_edits(
                api_base, api_key, model_name, safe_prompt, size_pixel,
                ref_files, quality=quality, n=1)
            self.log_request("POST", log_data["url"], log_data["payload"])
        else:
            img_data, log_data = openai_generations(
                api_base, api_key, model_name, safe_prompt, size_pixel,
                quality=quality, n=1)
            self.log_request("POST", log_data["url"], log_data["payload"])

        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        with open(filepath, "wb") as f:
            f.write(img_data)
        self.log_success(f"#{seq} 完成: {os.path.basename(filepath)} "
                         f"({len(img_data)/1024/1024:.2f} MB)")
        return filepath

    def _batch_done(self, ctx, ok, fail, total):
        self.generating = False
        ctx.progress_bar.set(1.0)
        # 恢复两个按钮
        self.ctx_single.btn_generate.configure(
            text="🚀  开始生成", state="normal",
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            text_color="white")
        self.ctx_batch.btn_generate.configure(
            text="🚀  批量生成", state="normal",
            fg_color=COLORS["accent"], hover_color=COLORS["accent_hover"],
            text_color="white")

        if fail == 0:
            self._set_status(f"✅ 完成：成功 {ok}/{total} 张", COLORS["success"])
        else:
            self._set_status(f"⚠️ 结束：成功 {ok}，失败 {fail}，共 {total} 张",
                             COLORS["warning"])
        self.log_info(f"结束 | 成功 {ok} | 失败 {fail} | 共 {total}")
        self.win.after(800, lambda: ctx.progress_frame.pack_forget())

    def _ui(self, fn):
        self.win.after(0, fn)

    def _open_file(self, path):
        try:
            if os.name == "nt":
                os.startfile(path)
            else:
                import subprocess
                subprocess.Popen(["xdg-open", path])
        except Exception as e:
            self.log_error(f"打开失败: {e}")

    def run(self):
        self.win.mainloop()


if __name__ == "__main__":
    app = ImageGeneratorApp()
    app.run()
